"""
Premium Consulting-Grade Presentation Generator
Modern management-consulting style with navy/blue/green accents
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from datetime import datetime

# Consulting Color Palette - Premium, Clean
NAVY = RGBColor(30, 50, 80)              # Deep navy primary
DEEP_NAVY = RGBColor(20, 35, 60)         # Darker navy for headers
BRIGHT_BLUE = RGBColor(60, 130, 200)     # Bright blue accent
SOFT_BLUE = RGBColor(180, 210, 240)      # Light blue panels
SOFT_GREEN = RGBColor(120, 180, 140)     # Soft green accent
LIGHT_GREEN = RGBColor(200, 235, 210)    # Light green panels
WHITE = RGBColor(255, 255, 255)          # Pure white
OFF_WHITE = RGBColor(250, 250, 252)      # Slightly warm white
LIGHT_GRAY = RGBColor(240, 242, 245)     # Soft panels
MED_GRAY = RGBColor(200, 205, 210)       # Borders/dividers
CHARCOAL = RGBColor(50, 55, 65)          # Primary text
SLATE = RGBColor(100, 110, 125)          # Secondary text
WARNING_RED = RGBColor(200, 80, 80)      # Red for warnings/assumptions

def add_consulting_header(slide, show_accent=True):
    """Clean consulting header with optional navy accent line"""
    if show_accent:
        # Thin navy accent line at top
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.04))
        accent.fill.solid()
        accent.fill.fore_color.rgb = NAVY
        accent.line.fill.background()

def add_consulting_footer(slide, page_num=""):
    """Minimal consulting footer"""
    # Separator line
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                  Inches(0.6), Inches(7.0), 
                                  Inches(8.8), Inches(0.01))
    sep.fill.solid()
    sep.fill.fore_color.rgb = MED_GRAY
    sep.line.fill.background()
    
    # Footer text
    footer = slide.shapes.add_textbox(Inches(0.6), Inches(7.08), Inches(8.8), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "XYZ Mobile App Diagnostic | Confidential | January 2026"
    p.font.size = Pt(8)
    p.font.color.rgb = SLATE
    p.font.name = "Calibri Light"
    
    if page_num:
        p2 = tf.add_paragraph()
        p2.text = page_num
        p2.font.size = Pt(8)
        p2.font.color.rgb = SLATE
        p2.alignment = PP_ALIGN.RIGHT

def add_slide_title_consulting(slide, title, top=0.35, font_size=26, color=NAVY):
    """Consulting-style slide title"""
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(8.8), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Calibri Light"
    return title_box

def add_subtitle_consulting(slide, subtitle, left=0.6, top=1.0, color=SLATE):
    """Professional subtitle"""
    sub_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(8.8), Inches(0.4))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(12)
    p.font.color.rgb = color
    p.font.name = "Calibri"
    return sub_box

def add_body_text_consulting(slide, bullets, left=0.6, top=1.4, width=8.8, font_size=11):
    """Clean consulting bullet text"""
    body = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = CHARCOAL
        p.font.name = "Calibri"
        p.space_after = Pt(8)
    return body

def add_card(slide, left, top, width, height, bg_color=LIGHT_GRAY, border_color=MED_GRAY):
    """Create a soft panel/card"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = border_color
    card.line.width = Pt(0.5)
    return card

def add_icon_circle(slide, left, top, size=0.35, color=BRIGHT_BLUE):
    """Add small icon circle placeholder"""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(left), Inches(top),
                                     Inches(size), Inches(size))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    return circle

def add_small_label(slide, text, left, top, color=NAVY, font_size=9):
    """Small label text"""
    label = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(2), Inches(0.25))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Calibri"
    return label

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ============ SLIDE 1 — CONSULTING COVER ============
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

# Clean white background
bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE
bg.line.fill.background()

# Top navy accent bar
navy_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.08))
navy_bar.fill.solid()
navy_bar.fill.fore_color.rgb = NAVY
navy_bar.line.fill.background()

# Client logo placeholder (top left)
logo_client = add_card(slide1, 0.6, 0.25, 1.2, 0.5, WHITE, MED_GRAY)
logo_text = slide1.shapes.add_textbox(Inches(0.7), Inches(0.32), Inches(1), Inches(0.4))
tf = logo_text.text_frame
p = tf.paragraphs[0]
p.text = "Client Logo"
p.font.size = Pt(9)
p.font.color.rgb = SLATE
p.alignment = PP_ALIGN.CENTER

# AB Brand placeholder (top right)
logo_ab = add_card(slide1, 8.2, 0.25, 1.2, 0.5, WHITE, MED_GRAY)
ab_text = slide1.shapes.add_textbox(Inches(8.3), Inches(0.32), Inches(1), Inches(0.4))
tf = ab_text.text_frame
p = tf.paragraphs[0]
p.text = "AB Brand"
p.font.size = Pt(9)
p.font.color.rgb = SLATE
p.alignment = PP_ALIGN.CENTER

# Left side - Title area (2/3 width)
title_area = slide1.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(6), Inches(3))
tf = title_area.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "XYZ Mobile App"
p.font.size = Pt(42)
p.font.bold = True
p.font.color.rgb = NAVY
p.font.name = "Calibri Light"

p = tf.add_paragraph()
p.text = "Performance & Release Management Diagnostic"
p.font.size = Pt(24)
p.font.color.rgb = CHARCOAL
p.font.name = "Calibri Light"
p.space_before = Pt(12)

p = tf.add_paragraph()
p.text = "XYZ Company"
p.font.size = Pt(16)
p.font.color.rgb = SLATE
p.font.name = "Calibri"
p.space_before = Pt(20)

p = tf.add_paragraph()
p.text = "January 2026"
p.font.size = Pt(14)
p.font.color.rgb = SLATE
p.font.name = "Calibri Light"
p.space_before = Pt(8)

# Right side - Abstract shapes suggesting analytics (1/3 width)
# Circle 1
c1 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), Inches(2.5), Inches(1.2), Inches(1.2))
c1.fill.solid()
c1.fill.fore_color.rgb = SOFT_BLUE
c1.line.fill.background()

# Circle 2
c2 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.4), Inches(3.2), Inches(0.8), Inches(0.8))
c2.fill.solid()
c2.fill.fore_color.rgb = BRIGHT_BLUE
c2.line.fill.background()

# Circle 3
c3 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.8), Inches(4.0), Inches(0.6), Inches(0.6))
c3.fill.solid()
c3.fill.fore_color.rgb = SOFT_GREEN
c3.line.fill.background()

# Small accent bars suggesting data
for i in range(4):
    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(7.0 + i * 0.4), Inches(4.8),
                                   Inches(0.25), Inches(0.3 + i * 0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY if i % 2 == 0 else BRIGHT_BLUE
    bar.line.fill.background()

# Bottom left footer
footer1 = slide1.shapes.add_textbox(Inches(0.6), Inches(6.8), Inches(4), Inches(0.3))
tf = footer1.text_frame
p = tf.paragraphs[0]
p.text = "XYZ Mobile App Diagnostic"
p.font.size = Pt(9)
p.font.color.rgb = SLATE

# Bottom separator
sep = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.0), Inches(8.8), Inches(0.01))
sep.fill.solid()
sep.fill.fore_color.rgb = MED_GRAY
sep.line.fill.background()

bottom_footer = slide1.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(8.8), Inches(0.3))
tf = bottom_footer.text_frame
p = tf.paragraphs[0]
p.text = "Confidential | Prepared for XYZ Company | January 2026"
p.font.size = Pt(8)
p.font.color.rgb = SLATE

# ============ SLIDE 2 — OUR UNDERSTANDING OF SCOPE ============
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_consulting_header(slide2)
add_consulting_footer(slide2, "02")

add_slide_title_consulting(slide2, "Our Understanding of Scope")
add_subtitle_consulting(slide2, "Client seeks a diagnostic-driven assessment to:")

# Six bullet cards in two columns
scope_items = [
    ("Analyze", "mobile app latency across Home, Insurance, Spend Track, Quiz & other flows"),
    ("Identify", "root causes behind long load times (6 seconds vs market 2–3 sec benchmark)"),
    ("Understand", "app size inflation (Android: 160MB → 400+MB installed; iOS: 402MB)"),
    ("Determine", "feasibility of moving to a monthly release cycle"),
    ("Recommend", "fixes backed by measurable RCA (no assumptions)"),
    ("Provide", "a North Star performance vision to guide long-term optimization")
]

y_start = 1.6
for i, (action, desc) in enumerate(scope_items):
    col = i % 2
    row = i // 2
    x = 0.6 if col == 0 else 5.3
    y = y_start + row * 1.15
    
    # Card background
    card = add_card(slide2, x, y, 4.3, 1.0, LIGHT_GRAY, MED_GRAY)
    
    # Icon circle
    icon = add_icon_circle(slide2, x + 0.15, y + 0.15, 0.3, BRIGHT_BLUE if i % 2 == 0 else SOFT_GREEN)
    
    # Action word (bold)
    action_box = slide2.shapes.add_textbox(Inches(x + 0.55), Inches(y + 0.18), Inches(0.8), Inches(0.3))
    tf = action_box.text_frame
    p = tf.paragraphs[0]
    p.text = action
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    # Description
    desc_box = slide2.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.55), Inches(4.0), Inches(0.4))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = CHARCOAL

# ============ SLIDE 3 — SCOPE OF DIAGNOSTIC ============
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_consulting_header(slide3)
add_consulting_footer(slide3, "03")

add_slide_title_consulting(slide3, "Scope of Diagnostic")
add_subtitle_consulting(slide3, "Mapped to Reference Structure")

# Three vertical cards
col_width = 2.9
x_positions = [0.6, 3.6, 6.6]

# Left card - AB Team
card1 = add_card(slide3, x_positions[0], 1.4, col_width, 4.8, OFF_WHITE, MED_GRAY)

# Header bar
header1 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                   Inches(x_positions[0]), Inches(1.4),
                                   Inches(col_width), Inches(0.5))
header1.fill.solid()
header1.fill.fore_color.rgb = NAVY
header1.line.fill.background()

header1_text = slide3.shapes.add_textbox(Inches(x_positions[0]), Inches(1.48), Inches(col_width), Inches(0.4))
tf = header1_text.text_frame
p = tf.paragraphs[0]
p.text = "Project Start – Pre‑Requisite"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# AB Team label
team_label = add_small_label(slide3, "AB Team", x_positions[0] + 0.15, 2.05, NAVY, 10)

# Team list
team_items = [
    "Mobile Performance Lead (Flutter)",
    "Mobile Performance Engineer", 
    "Engineering Manager"
]
y = 2.35
for item in team_items:
    item_box = slide3.shapes.add_textbox(Inches(x_positions[0] + 0.25), Inches(y), Inches(2.5), Inches(0.35))
    tf = item_box.text_frame
    p = tf.paragraphs[0]
    p.text = "• " + item
    p.font.size = Pt(9)
    p.font.color.rgb = CHARCOAL
    y += 0.35

# Activities label
act_label = add_small_label(slide3, "Activities", x_positions[0] + 0.15, 3.55, NAVY, 10)

activities = [
    "Access setup: source code, UAT builds",
    "Journey & technical walkthrough",
    "Environment & build readiness confirmation"
]
y = 3.85
for item in activities:
    item_box = slide3.shapes.add_textbox(Inches(x_positions[0] + 0.25), Inches(y), Inches(2.5), Inches(0.35))
    tf = item_box.text_frame
    p = tf.paragraphs[0]
    p.text = "• " + item
    p.font.size = Pt(9)
    p.font.color.rgb = CHARCOAL
    y += 0.35

# Middle card - Client Inputs
card2 = add_card(slide3, x_positions[1], 1.4, col_width, 4.8, SOFT_BLUE, MED_GRAY)

header2 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x_positions[1]), Inches(1.4),
                                   Inches(col_width), Inches(0.5))
header2.fill.solid()
header2.fill.fore_color.rgb = BRIGHT_BLUE
header2.line.fill.background()

header2_text = slide3.shapes.add_textbox(Inches(x_positions[1]), Inches(1.48), Inches(col_width), Inches(0.4))
tf = header2_text.text_frame
p = tf.paragraphs[0]
p.text = "Diagnostic Pre‑Requisite"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Inputs label
inputs_label = add_small_label(slide3, "Inputs Needed From Client", x_positions[1] + 0.15, 2.05, BRIGHT_BLUE, 10)

client_inputs = [
    "Latest production build (APK/IPA)",
    "Access to Analytics, CMS",
    "API documentation",
    "Release pipeline documentation",
    "Third‑party SDK list"
]
y = 2.35
for item in client_inputs:
    item_box = slide3.shapes.add_textbox(Inches(x_positions[1] + 0.25), Inches(y), Inches(2.5), Inches(0.35))
    tf = item_box.text_frame
    p = tf.paragraphs[0]
    p.text = "• " + item
    p.font.size = Pt(9)
    p.font.color.rgb = CHARCOAL
    y += 0.35

# Right card - Placeholder
card3 = add_card(slide3, x_positions[2], 1.4, col_width, 4.8, LIGHT_GREEN, MED_GRAY)

header3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x_positions[2]), Inches(1.4),
                                   Inches(col_width), Inches(0.5))
header3.fill.solid()
header3.fill.fore_color.rgb = SOFT_GREEN
header3.line.fill.background()

header3_text = slide3.shapes.add_textbox(Inches(x_positions[2]), Inches(1.48), Inches(col_width), Inches(0.4))
tf = header3_text.text_frame
p = tf.paragraphs[0]
p.text = "Client Responsibilities"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

placeholder = slide3.shapes.add_textbox(Inches(x_positions[2] + 0.2), Inches(3), Inches(2.5), Inches(1))
tf = placeholder.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "[Placeholder for client-specific responsibilities and commitments]"
p.font.size = Pt(9)
p.font.italic = True
p.font.color.rgb = SLATE
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 4 — NORTH STAR VISION ============
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_consulting_header(slide4)
add_consulting_footer(slide4, "04")

add_slide_title_consulting(slide4, "North Star Vision")
add_subtitle_consulting(slide4, "Performance North Star (Industry Benchmarks)")

# Top comparison bar for Load Time
comp_bar_bg = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(0.6), Inches(1.5),
                                       Inches(8.8), Inches(1.0))
comp_bar_bg.fill.solid()
comp_bar_bg.fill.fore_color.rgb = LIGHT_GRAY
comp_bar_bg.line.color.rgb = MED_GRAY
comp_bar_bg.line.width = Pt(1)

# Title for comparison
load_title = slide4.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(4), Inches(0.3))
tf = load_title.text_frame
p = tf.paragraphs[0]
p.text = "Primary Screen Load Time"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = NAVY

# Current (red) portion
current_bar = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0.8), Inches(2.0),
                                       Inches(3.5), Inches(0.35))
current_bar.fill.solid()
current_bar.fill.fore_color.rgb = WARNING_RED
current_bar.line.fill.background()

current_text = slide4.shapes.add_textbox(Inches(0.8), Inches(2.05), Inches(3.5), Inches(0.3))
tf = current_text.text_frame
p = tf.paragraphs[0]
p.text = "Current: 6.0 sec (Slow)"
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Target (green) portion
target_bar = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(4.4), Inches(2.0),
                                      Inches(3.5), Inches(0.35))
target_bar.fill.solid()
target_bar.fill.fore_color.rgb = SOFT_GREEN
target_bar.line.fill.background()

target_text = slide4.shapes.add_textbox(Inches(4.4), Inches(2.05), Inches(3.5), Inches(0.3))
tf = target_text.text_frame
p = tf.paragraphs[0]
p.text = "Market Benchmark: 2.0 sec"
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Arrow showing improvement
arrow = slide4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(7.9), Inches(2.05),
                                 Inches(1.3), Inches(0.25))
arrow.fill.solid()
arrow.fill.fore_color.rgb = NAVY
arrow.line.fill.background()

improvement = slide4.shapes.add_textbox(Inches(7.9), Inches(1.75), Inches(1.3), Inches(0.25))
tf = improvement.text_frame
p = tf.paragraphs[0]
p.text = "67% faster"
p.font.size = Pt(9)
p.font.bold = True
p.font.color.rgb = SOFT_GREEN
p.alignment = PP_ALIGN.CENTER

# Six benchmark cards in 2x3 grid
benchmarks = [
    ("Tab-switch latency", "150–250 ms"),
    ("App size target", "30–40% reduction"),
    ("API latency goal", "<150 ms critical"),
    ("Frame stability", "<16 ms/frame"),
    ("Release cadence", "Monthly train"),
    ("Current trajectory", "Action required")
]

positions = [(0.6, 2.8), (3.5, 2.8), (6.4, 2.8), (0.6, 4.3), (3.5, 4.3), (6.4, 4.3)]
for i, (label, value) in enumerate(benchmarks):
    x, y = positions[i]
    
    card = add_card(slide4, x, y, 2.7, 1.3, OFF_WHITE, MED_GRAY)
    
    # Icon
    icon = add_icon_circle(slide4, x + 0.15, y + 0.15, 0.25, 
                           BRIGHT_BLUE if i < 3 else SOFT_GREEN)
    
    # Label
    label_box = slide4.shapes.add_textbox(Inches(x + 0.5), Inches(y + 0.18), Inches(2.0), Inches(0.3))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    # Value
    value_box = slide4.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.55), Inches(2.4), Inches(0.5))
    tf = value_box.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CHARCOAL

# Note bar at bottom
note_bar = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(0.6), Inches(6.1),
                                    Inches(8.8), Inches(0.5))
note_bar.fill.solid()
note_bar.fill.fore_color.rgb = RGBColor(255, 250, 230)
note_bar.line.color.rgb = RGBColor(230, 200, 120)
note_bar.line.width = Pt(1)

note_text = slide4.shapes.add_textbox(Inches(0.8), Inches(6.25), Inches(8.4), Inches(0.3))
tf = note_text.text_frame
p = tf.paragraphs[0]
p.text = "⚠ These are reference benchmarks only. Final commitments will be established after RCA is completed."
p.font.size = Pt(10)
p.font.color.rgb = CHARCOAL

# ============ SLIDE 5 — ASSUMPTIONS ============
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_consulting_header(slide5)
add_consulting_footer(slide5, "05")

add_slide_title_consulting(slide5, "Assumptions")

# Left panel - Critical Success Factors
csf_card = add_card(slide5, 0.6, 1.3, 3.0, 4.5, SOFT_BLUE, MED_GRAY)

csf_icon = add_icon_circle(slide5, 0.9, 1.6, 0.4, BRIGHT_BLUE)

csf_title = slide5.shapes.add_textbox(Inches(1.4), Inches(1.65), Inches(2.2), Inches(0.4))
tf = csf_title.text_frame
p = tf.paragraphs[0]
p.text = "Critical Success Factors"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = NAVY

csf_desc = slide5.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(2.4), Inches(1.5))
tf = csf_desc.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "These assumptions guide the diagnostic approach and determine feasibility outcomes."
p.font.size = Pt(10)
p.font.color.rgb = CHARCOAL

# Right side - assumption strips
assumptions_data = [
    ("All access (code, builds, dashboards) will be provided by the client", "Prerequisite", SOFT_BLUE),
    ("Third‑party SDK behavior and CMS limitations may restrict optimization", "Technical Constraint", RGBColor(255, 235, 200)),
    ("No changes to backend or CMS unless explicitly included", "Scope Boundary", LIGHT_GREEN),
    ("RCA outcomes will determine feasibility of performance enhancements", "Methodology", SOFT_BLUE),
    ("Recommendations will be measurable and derived from profiling & data", "Data-Driven", LIGHT_GREEN),
    ("Business‑driven UI/UX changes are out of scope unless mutually agreed", "Scope Boundary", RGBColor(255, 235, 200)),
    ("Release Management changes are advisory; implementation requires client DevOps", "Advisory Role", SOFT_BLUE)
]

y = 1.3
for text, label, color in assumptions_data:
    # Strip background
    strip = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(3.9), Inches(y),
                                     Inches(5.5), Inches(0.58))
    strip.fill.solid()
    strip.fill.fore_color.rgb = color
    strip.line.color.rgb = MED_GRAY
    strip.line.width = Pt(0.5)
    
    # Label
    label_box = slide5.shapes.add_textbox(Inches(4.0), Inches(y + 0.05), Inches(1.5), Inches(0.2))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = BRIGHT_BLUE if color == SOFT_BLUE else SOFT_GREEN if color == LIGHT_GREEN else RGBColor(200, 150, 50)
    
    # Text
    text_box = slide5.shapes.add_textbox(Inches(4.0), Inches(y + 0.22), Inches(5.2), Inches(0.35))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = CHARCOAL
    
    y += 0.68

# ============ SLIDE 6 — ARCHITECTURE ============
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_consulting_header(slide6)
add_consulting_footer(slide6, "06")

add_slide_title_consulting(slide6, "Architecture & Design Considerations")
add_subtitle_consulting(slide6, "Evaluating the mobile application ecosystem to support performance goals.")

# Two columns
# Left - Core Architecture
left_card = add_card(slide6, 0.6, 1.5, 4.2, 4.5, OFF_WHITE, MED_GRAY)

left_header = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0.6), Inches(1.5),
                                       Inches(4.2), Inches(0.5))
left_header.fill.solid()
left_header.fill.fore_color.rgb = NAVY
left_header.line.fill.background()

left_title = slide6.shapes.add_textbox(Inches(0.6), Inches(1.58), Inches(4.2), Inches(0.4))
tf = left_title.text_frame
p = tf.paragraphs[0]
p.text = "Core Architecture"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Icon
left_icon = add_icon_circle(slide6, 2.4, 2.15, 0.3, BRIGHT_BLUE)

left_items = [
    "Modular, layered architecture assessment",
    "Asynchronous vs synchronous rendering optimization",
    "Lazy‑loading feasibility",
    "Separation of concerns for future scalability"
]
y = 2.6
for item in left_items:
    item_box = slide6.shapes.add_textbox(Inches(0.9), Inches(y), Inches(3.6), Inches(0.4))
    tf = item_box.text_frame
    p = tf.paragraphs[0]
    p.text = "• " + item
    p.font.size = Pt(10)
    p.font.color.rgb = CHARCOAL
    y += 0.5

# Right - Connectivity & Governance
right_card = add_card(slide6, 5.2, 1.5, 4.2, 4.5, LIGHT_GREEN, MED_GRAY)

right_header = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(5.2), Inches(1.5),
                                        Inches(4.2), Inches(0.5))
right_header.fill.solid()
right_header.fill.fore_color.rgb = SOFT_GREEN
right_header.line.fill.background()

right_title = slide6.shapes.add_textbox(Inches(5.2), Inches(1.58), Inches(4.2), Inches(0.4))
tf = right_title.text_frame
p = tf.paragraphs[0]
p.text = "Connectivity & Governance"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

right_icon = add_icon_circle(slide6, 7.0, 2.15, 0.3, SOFT_GREEN)

right_items = [
    "API sequencing, dependency mapping",
    "Third‑party SDK footprint & load behavior",
    "Asset compression & caching strategies",
    "Release governance & branching strategy review"
]
y = 2.6
for item in right_items:
    item_box = slide6.shapes.add_textbox(Inches(5.5), Inches(y), Inches(3.6), Inches(0.4))
    tf = item_box.text_frame
    p = tf.paragraphs[0]
    p.text = "• " + item
    p.font.size = Pt(10)
    p.font.color.rgb = CHARCOAL
    y += 0.5

# ============ SLIDE 7 — DIAGNOSTIC ARCHITECTURE VIEW ============
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_consulting_header(slide7)
add_consulting_footer(slide7, "07")

add_slide_title_consulting(slide7, "Proposed Diagnostic Architecture View")

# Architecture diagram boxes
# Left - App Frontend
app_box = add_card(slide7, 0.5, 1.5, 2.5, 2.5, SOFT_BLUE, BRIGHT_BLUE)
app_title = slide7.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(2.1), Inches(0.3))
tf = app_title.text_frame
p = tf.paragraphs[0]
p.text = "App Frontend"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = NAVY

app_sub = slide7.shapes.add_textbox(Inches(0.7), Inches(1.95), Inches(2.1), Inches(0.25))
tf = app_sub.text_frame
p = tf.paragraphs[0]
p.text = "(Flutter)"
p.font.size = Pt(9)
p.font.color.rgb = SLATE

app_items = ["UI", "State Mgmt", "Local Cache"]
y = 2.4
for item in app_items:
    item_box = slide7.shapes.add_textbox(Inches(0.7), Inches(y), Inches(2.1), Inches(0.3))
    tf = item_box.text_frame
    p = tf.paragraphs[0]
    p.text = "• " + item
    p.font.size = Pt(9)
    p.font.color.rgb = CHARCOAL
    y += 0.3

# Center - API Gateway
gateway_box = add_card(slide7, 3.5, 1.8, 3, 1.8, LIGHT_GRAY, MED_GRAY)
gateway_title = slide7.shapes.add_textbox(Inches(3.7), Inches(1.95), Inches(2.6), Inches(0.3))
tf = gateway_title.text_frame
p = tf.paragraphs[0]
p.text = "Interaction Layer"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = NAVY

gateway_items = ["API Gateway", "CMS Modules"]
y = 2.35
for item in gateway_items:
    item_box = slide7.shapes.add_textbox(Inches(3.7), Inches(y), Inches(2.6), Inches(0.3))
    tf = item_box.text_frame
    p = tf.paragraphs[0]
    p.text = "• " + item
    p.font.size = Pt(9)
    p.font.color.rgb = CHARCOAL
    y += 0.3

# Right - Backend
backend_box = add_card(slide7, 7, 1.5, 2.5, 2.5, OFF_WHITE, MED_GRAY)
backend_title = slide7.shapes.add_textbox(Inches(7.2), Inches(1.6), Inches(2.1), Inches(0.3))
tf = backend_title.text_frame
p = tf.paragraphs[0]
p.text = "Backend Systems"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = NAVY

backend_sub = slide7.shapes.add_textbox(Inches(7.2), Inches(1.95), Inches(2.1), Inches(0.25))
tf = backend_sub.text_frame
p = tf.paragraphs[0]
p.text = "(Core Services)"
p.font.size = Pt(9)
p.font.color.rgb = SLATE

# SDKs below frontend
sdk_box = add_card(slide7, 0.5, 4.3, 2.5, 1.5, LIGHT_GREEN, SOFT_GREEN)
sdk_title = slide7.shapes.add_textbox(Inches(0.7), Inches(4.4), Inches(2.1), Inches(0.3))
tf = sdk_title.text_frame
p = tf.paragraphs[0]
p.text = "Third‑Party SDKs"
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = NAVY

sdk_items = ["Fly", "MarTech", "Payments", "Firebase"]
x = 0.7
for item in sdk_items:
    item_box = slide7.shapes.add_textbox(Inches(x), Inches(4.9), Inches(0.6), Inches(0.3))
    tf = item_box.text_frame
    p = tf.paragraphs[0]
    p.text = item
    p.font.size = Pt(8)
    p.font.color.rgb = CHARCOAL
    x += 0.55

# Bottom - Telemetry bar
telem_bar = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0.5), Inches(6.0),
                                     Inches(9), Inches(0.6))
telem_bar.fill.solid()
telem_bar.fill.fore_color.rgb = DEEP_NAVY
telem_bar.line.fill.background()

telem_title = slide7.shapes.add_textbox(Inches(0.7), Inches(6.15), Inches(2), Inches(0.3))
tf = telem_title.text_frame
p = tf.paragraphs[0]
p.text = "Performance Telemetry & Analytics"
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = WHITE

# Telemetry tags
tags = ["App Start", "API Latency", "Memory", "Crash/ANR"]
x = 3.5
for tag in tags:
    tag_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x), Inches(6.15),
                                       Inches(1.3), Inches(0.3))
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = BRIGHT_BLUE
    tag_box.line.fill.background()
    
    tag_text = slide7.shapes.add_textbox(Inches(x), Inches(6.2), Inches(1.3), Inches(0.25))
    tf = tag_text.text_frame
    p = tf.paragraphs[0]
    p.text = tag
    p.font.size = Pt(8)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    x += 1.4

# Arrows (simplified representation)
arrow1 = slide7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.05), Inches(2.4), Inches(0.4), Inches(0.2))
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = MED_GRAY

arrow2 = slide7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.55), Inches(2.4), Inches(0.4), Inches(0.2))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = MED_GRAY

# ============ SLIDE 8 — APPROACH TIMELINE ============
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_consulting_header(slide8)
add_consulting_footer(slide8, "08")

add_slide_title_consulting(slide8, "Our Approach – Diagnostic")
add_subtitle_consulting(slide8, "Week 0 to Week 4 Execution Plan")

# Timeline steps
phases = [
    ("Week 0", "Setup & Access", ["Access provisioning", "Environment setup", "Build validation"], NAVY),
    ("Week 1", "Profiling & Analysis", ["Load‑time benchmarking", "API call mapping", "Size breakdown", "Cache profiling"], BRIGHT_BLUE),
    ("Week 2", "RCA & Observations", ["Root‑cause identification", "Fixable analysis", "SDK constraints", "Pipeline review"], SOFT_GREEN),
    ("Week 3", "Discussions & Verifications", ["Joint walkthrough", "Client validation", "Feasibility check", "Release alignment"], NAVY),
    ("Week 4", "Final Report", ["Target alignment", "Uplift range", "Recommended fixes", "Execution plan"], BRIGHT_BLUE)
]

x_start = 0.5
step_width = 1.8
for i, (week, phase, items, color) in enumerate(phases):
    x = x_start + (i * step_width)
    
    # Step box
    step = add_card(slide8, x, 1.6, 1.7, 4.5, OFF_WHITE, MED_GRAY)
    
    # Week header
    week_header = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           Inches(x), Inches(1.6),
                                           Inches(1.7), Inches(0.5))
    week_header.fill.solid()
    week_header.fill.fore_color.rgb = color
    week_header.line.fill.background()
    
    week_text = slide8.shapes.add_textbox(Inches(x), Inches(1.68), Inches(1.7), Inches(0.4))
    tf = week_text.text_frame
    p = tf.paragraphs[0]
    p.text = week
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Phase name
    phase_box = slide8.shapes.add_textbox(Inches(x + 0.1), Inches(2.2), Inches(1.5), Inches(0.6))
    tf = phase_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = phase
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    
    # Items
    y = 2.95
    for item in items:
        item_box = slide8.shapes.add_textbox(Inches(x + 0.15), Inches(y), Inches(1.5), Inches(0.5))
        tf = item_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "• " + item
        p.font.size = Pt(8)
        p.font.color.rgb = CHARCOAL
        y += 0.55

# Connector line
connector = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(1.2), Inches(3.85),
                                     Inches(7.6), Inches(0.02))
connector.fill.solid()
connector.fill.fore_color.rgb = MED_GRAY
connector.line.fill.background()

# ============ SLIDE 9 — TEAM STRUCTURE ============
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_consulting_header(slide9)
add_consulting_footer(slide9, "09")

add_slide_title_consulting(slide9, "Team Structure – Diagnostic Phase")

# Table header
headers = ["Role", "Count", "Responsibility"]
x_positions = [0.6, 4.0, 4.8]
col_widths = [3.3, 0.7, 4.5]

for i, (header, x, w) in enumerate(zip(headers, x_positions, col_widths)):
    cell = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x), Inches(1.5),
                                    Inches(w), Inches(0.45))
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    cell.line.fill.background()
    
    text = slide9.shapes.add_textbox(Inches(x), Inches(1.58), Inches(w), Inches(0.35))
    tf = text.text_frame
    p = tf.paragraphs[0]
    p.text = header
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# Team rows
team = [
    ("Mobile Performance Lead", "1", "Profiling, rendering, app load optimization"),
    ("Mobile Engineer", "2", "Code analysis, architecture assessment"),
    ("Solution Architect", "1", "Architecture oversight and technical guidance"),
    ("Engineering Manager / Program Management", "1", "Oversight, coordination, stakeholder management")
]

y_pos = 2.05
for role, count, resp in team:
    # Alternating row colors
    bg_color = OFF_WHITE if y_pos < 3.5 else LIGHT_GRAY
    
    # Role cell
    role_bg = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0.6), Inches(y_pos),
                                       Inches(3.3), Inches(0.55))
    role_bg.fill.solid()
    role_bg.fill.fore_color.rgb = bg_color
    role_bg.line.color.rgb = MED_GRAY
    role_bg.line.width = Pt(0.25)
    
    role_text = slide9.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.15), Inches(3.0), Inches(0.3))
    tf = role_text.text_frame
    p = tf.paragraphs[0]
    p.text = role
    p.font.size = Pt(10)
    p.font.color.rgb = CHARCOAL
    
    # Count cell
    count_bg = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(4.0), Inches(y_pos),
                                        Inches(0.7), Inches(0.55))
    count_bg.fill.solid()
    count_bg.fill.fore_color.rgb = SOFT_BLUE
    count_bg.line.color.rgb = MED_GRAY
    count_bg.line.width = Pt(0.25)
    
    count_text = slide9.shapes.add_textbox(Inches(4.0), Inches(y_pos + 0.15), Inches(0.7), Inches(0.3))
    tf = count_text.text_frame
    p = tf.paragraphs[0]
    p.text = count
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    
    # Responsibility cell
    resp_bg = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(4.8), Inches(y_pos),
                                       Inches(4.5), Inches(0.55))
    resp_bg.fill.solid()
    resp_bg.fill.fore_color.rgb = bg_color
    resp_bg.line.color.rgb = MED_GRAY
    resp_bg.line.width = Pt(0.25)
    
    resp_text = slide9.shapes.add_textbox(Inches(5.0), Inches(y_pos + 0.1), Inches(4.1), Inches(0.4))
    tf = resp_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = resp
    p.font.size = Pt(9)
    p.font.color.rgb = CHARCOAL
    
    y_pos += 0.6

# Total badge
badge = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(7.5), Inches(5.0),
                                 Inches(1.8), Inches(0.5))
badge.fill.solid()
badge.fill.fore_color.rgb = SOFT_GREEN
badge.line.fill.background()

badge_text = slide9.shapes.add_textbox(Inches(7.5), Inches(5.12), Inches(1.8), Inches(0.35))
tf = badge_text.text_frame
p = tf.paragraphs[0]
p.text = "Total: 3 Members"
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 10 — GANTT TIMELINE ============
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_consulting_header(slide10)
add_consulting_footer(slide10, "10")

add_slide_title_consulting(slide10, "Gantt Timeline – Diagnostic Phase")
add_subtitle_consulting(slide10, "Weeks 0–4")

# Week headers
weeks = ["Week 0", "Week 1", "Week 2", "Week 3", "Week 4"]
for i, week in enumerate(weeks):
    header = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(3.5 + i * 1.2), Inches(1.5),
                                       Inches(1.15), Inches(0.4))
    header.fill.solid()
    header.fill.fore_color.rgb = OFF_WHITE
    header.line.color.rgb = MED_GRAY
    header.line.width = Pt(0.5)
    
    text = slide10.shapes.add_textbox(Inches(3.5 + i * 1.2), Inches(1.58), Inches(1.15), Inches(0.3))
    tf = text.text_frame
    p = tf.paragraphs[0]
    p.text = week
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

# Gantt bars
tasks = [
    ("Setup & Access", 0, 1, NAVY),
    ("Profiling & Analysis", 1, 1, BRIGHT_BLUE),
    ("RCA & Observations", 2, 1, SOFT_GREEN),
    ("Discussions & Verifications", 3, 1, NAVY),
    ("Final Report", 4, 1, BRIGHT_BLUE)
]

y_pos = 2.2
for task, start, duration, color in tasks:
    # Task label
    label = slide10.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.1), Inches(2.8), Inches(0.35))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = task
    p.font.size = Pt(10)
    p.font.color.rgb = CHARCOAL
    
    # Bar
    bar = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(3.5 + start * 1.2), Inches(y_pos),
                                    Inches(duration * 1.1), Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    
    y_pos += 0.55

# Milestone markers
milestones = [4.1, 5.3, 6.5, 7.7, 8.9]
for i, x in enumerate(milestones):
    marker = slide10.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(x - 0.06), Inches(2.05),
                                       Inches(0.12), Inches(0.12))
    marker.fill.solid()
    marker.fill.fore_color.rgb = WHITE
    marker.line.color.rgb = MED_GRAY
    marker.line.width = Pt(2)

# ============ SLIDE 11 — COMMERCIAL STRUCTURE ============
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_consulting_header(slide11)
add_consulting_footer(slide11, "11")

add_slide_title_consulting(slide11, "Commercial Structure – Template")

# Diagnostic Phase panel
diag_panel = add_card(slide11, 0.6, 1.5, 8.8, 2.2, OFF_WHITE, MED_GRAY)

# Header strip
diag_header = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.6), Inches(1.5),
                                        Inches(8.8), Inches(0.5))
diag_header.fill.solid()
diag_header.fill.fore_color.rgb = NAVY
diag_header.line.fill.background()

diag_title = slide11.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(4), Inches(0.35))
tf = diag_title.text_frame
p = tf.paragraphs[0]
p.text = "🔍 Diagnostic Phase (4 Weeks)"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = WHITE

diag_items = [
    "Fixed fee (based on 4–5 resources for 1 month equivalent)",
    "Covers analysis, RCA, reporting, release advisory"
]
y = 2.15
for item in diag_items:
    item_box = slide11.shapes.add_textbox(Inches(0.9), Inches(y), Inches(8), Inches(0.35))
    tf = item_box.text_frame
    p = tf.paragraphs[0]
    p.text = "• " + item
    p.font.size = Pt(10)
    p.font.color.rgb = CHARCOAL
    y += 0.35

# Execution Phase panel
exec_panel = add_card(slide11, 0.6, 4.0, 8.8, 2.0, LIGHT_GRAY, MED_GRAY)

exec_header = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.6), Inches(4.0),
                                        Inches(8.8), Inches(0.5))
exec_header.fill.solid()
exec_header.fill.fore_color.rgb = BRIGHT_BLUE
exec_header.line.fill.background()

exec_title = slide11.shapes.add_textbox(Inches(0.9), Inches(4.1), Inches(3), Inches(0.35))
tf = exec_title.text_frame
p = tf.paragraphs[0]
p.text = "🚀 Execution Phase"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = WHITE

exec_items = [
    "To be estimated based on diagnostic output",
    "Dependent on size of fixable items"
]
y = 4.65
for item in exec_items:
    item_box = slide11.shapes.add_textbox(Inches(0.9), Inches(y), Inches(8), Inches(0.35))
    tf = item_box.text_frame
    p = tf.paragraphs[0]
    p.text = "• " + item
    p.font.size = Pt(10)
    p.font.color.rgb = CHARCOAL
    y += 0.35

# ============ SLIDE 12 — RISKS & DEPENDENCIES ============
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_consulting_header(slide12)
add_consulting_footer(slide12, "12")

add_slide_title_consulting(slide12, "Risks & Dependencies")

# Risk grid
risks = [
    ("Third‑party SDK limitations", "External dependencies"),
    ("CMS payload constraints", "Platform constraints"),
    ("Launch‑time API dependencies", "Performance blocker"),
    ("Device fragmentation & low‑RAM behavior", "Compatibility risk"),
    ("Release process maturity", "Operational readiness"),
    ("Environment availability", "Access & provisioning")
]

positions = [(0.6, 1.5), (5.3, 1.5), (0.6, 3.0), (5.3, 3.0), (0.6, 4.5), (5.3, 4.5)]
for i, (risk, category) in enumerate(risks):
    x, y = positions[i]
    
    # Card
    card = add_card(slide12, x, y, 4.2, 1.3, OFF_WHITE, MED_GRAY)
    
    # Warning icon
    warning = slide12.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(x + 0.15), Inches(y + 0.15),
                                        Inches(0.3), Inches(0.3))
    warning.fill.solid()
    warning.fill.fore_color.rgb = RGBColor(240, 180, 80)
    warning.line.fill.background()
    
    warning_text = slide12.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.2), Inches(0.3), Inches(0.25))
    tf = warning_text.text_frame
    p = tf.paragraphs[0]
    p.text = "!"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Risk text
    risk_box = slide12.shapes.add_textbox(Inches(x + 0.55), Inches(y + 0.2), Inches(3.5), Inches(0.6))
    tf = risk_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = risk
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = CHARCOAL
    
    # Category tag
    cat_box = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(x + 0.55), Inches(y + 0.85),
                                        Inches(1.5), Inches(0.25))
    cat_box.fill.solid()
    cat_box.fill.fore_color.rgb = LIGHT_GRAY
    cat_box.line.fill.background()
    
    cat_text = slide12.shapes.add_textbox(Inches(x + 0.55), Inches(y + 0.88), Inches(1.5), Inches(0.22))
    tf = cat_text.text_frame
    p = tf.paragraphs[0]
    p.text = category
    p.font.size = Pt(8)
    p.font.color.rgb = SLATE
    p.alignment = PP_ALIGN.CENTER

# Legend
legend = slide12.shapes.add_textbox(Inches(0.6), Inches(6.2), Inches(4), Inches(0.3))
tf = legend.text_frame
p = tf.paragraphs[0]
p.text = "Risk Levels: ⚠ High  ⚡ Medium  ✓ Low (to be assessed during diagnostic)"
p.font.size = Pt(9)
p.font.color.rgb = SLATE

# ============ SLIDE 13 — FINAL OUTCOME ============
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
add_consulting_header(slide13)
add_consulting_footer(slide13, "13")

add_slide_title_consulting(slide13, "Final Outcome")
add_subtitle_consulting(slide13, "North‑Star aligned, data‑backed performance roadmap")

# Four outcome cards
outcomes = [
    ("What can be improved", "Detailed technical recommendations with estimated impact", BRIGHT_BLUE),
    ("What cannot be improved", "Constraints documentation with technical rationale", WARNING_RED),
    ("Expected uplift range", "Quantified performance gains per optimization category", SOFT_GREEN),
    ("Team & timeline", "Recommended structure and implementation roadmap", NAVY)
]

x_positions = [0.6, 2.9, 5.2, 7.5]
for i, (title, desc, color) in enumerate(outcomes):
    x = x_positions[i]
    
    # Card
    card = add_card(slide13, x, 1.6, 2.1, 3.0, OFF_WHITE, MED_GRAY)
    
    # Color header
    header = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(x), Inches(1.6),
                                       Inches(2.1), Inches(0.4))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()
    
    # Icon placeholder
    icon = slide13.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(x + 0.85), Inches(2.15),
                                     Inches(0.4), Inches(0.4))
    icon.fill.solid()
    icon.fill.fore_color.rgb = WHITE
    icon.line.color.rgb = color
    icon.line.width = Pt(2)
    
    # Title
    title_box = slide13.shapes.add_textbox(Inches(x + 0.15), Inches(2.7), Inches(1.8), Inches(0.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    
    # Description
    desc_box = slide13.shapes.add_textbox(Inches(x + 0.15), Inches(3.4), Inches(1.8), Inches(1.0))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(9)
    p.font.color.rgb = CHARCOAL
    p.alignment = PP_ALIGN.CENTER

# Fifth highlight element
highlight = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.6), Inches(5.0),
                                      Inches(8.8), Inches(0.8))
highlight.fill.solid()
highlight.fill.fore_color.rgb = LIGHT_GREEN
highlight.line.color.rgb = SOFT_GREEN
highlight.line.width = Pt(2)

highlight_icon = slide13.shapes.add_shape(MSO_SHAPE.OVAL,
                                           Inches(1.0), Inches(5.25),
                                           Inches(0.4), Inches(0.4))
highlight_icon.fill.solid()
highlight_icon.fill.fore_color.rgb = WHITE
highlight_icon.line.color.rgb = SOFT_GREEN
highlight_icon.line.width = Pt(2)

highlight_title = slide13.shapes.add_textbox(Inches(1.6), Inches(5.25), Inches(3), Inches(0.35))
tf = highlight_title.text_frame
p = tf.paragraphs[0]
p.text = "Monthly Release Readiness Assessment"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = NAVY

highlight_desc = slide13.shapes.add_textbox(Inches(1.6), Inches(5.55), Inches(7.5), Inches(0.25))
tf = highlight_desc.text_frame
p = tf.paragraphs[0]
p.text = "Specific capability gaps and transition roadmap for predictable monthly releases"
p.font.size = Pt(10)
p.font.color.rgb = CHARCOAL

# Save with timestamp
timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
output_name = f'XYZ_Mobile_App_Diagnostic_Consulting_{timestamp}.pptx'
output_path = f'/Users/Mudassar.Hakim/tempfiles/{output_name}'

prs.save(output_path)
print(f"✓ Premium consulting presentation created: {output_name}")
print(f"  Style: Management-consulting, navy/blue/green, clean minimal design")
