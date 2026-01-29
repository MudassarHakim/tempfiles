"""
Creative Presentation Generator - XYZ Mobile App Diagnostic
Enhanced visuals, infographics, and engaging layouts
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from datetime import datetime

# Vibrant Color Palette
PRIMARY_RED = RGBColor(220, 38, 38)
DEEP_RED = RGBColor(153, 27, 27)
BRIGHT_YELLOW = RGBColor(251, 191, 36)
GOLD = RGBColor(245, 158, 11)
AMBER = RGBColor(217, 119, 6)
CREAM = RGBColor(254, 243, 199)
LIGHT_YELLOW = RGBColor(254, 252, 232)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(31, 41, 55)
DARK_GRAY = RGBColor(55, 65, 81)
ORANGE = RGBColor(249, 115, 22)
CORAL = RGBColor(251, 146, 60)

def add_gradient_header(slide, height=1.2):
    """Add gradient-style header with accent"""
    # Main header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(height))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_RED
    header.line.fill.background()
    
    # Yellow accent stripe at bottom of header
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(height - 0.1), Inches(10), Inches(0.1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BRIGHT_YELLOW
    accent.line.fill.background()
    
    # Decorative circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(0.2), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BRIGHT_YELLOW
    circle.line.fill.background()
    return header

def add_decorative_footer(slide):
    """Add decorative footer with wave pattern"""
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.0), Inches(10), Inches(0.5))
    footer.fill.solid()
    footer.fill.fore_color.rgb = DEEP_RED
    footer.line.fill.background()
    
    # Yellow line above footer
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.95), Inches(10), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = BRIGHT_YELLOW
    line.line.fill.background()

def add_title_with_icon(slide, title, icon_shape=MSO_SHAPE.OVAL, top=0.25):
    """Add title with decorative icon"""
    # Icon
    icon = slide.shapes.add_shape(icon_shape, Inches(0.5), Inches(top), Inches(0.5), Inches(0.5))
    icon.fill.solid()
    icon.fill.fore_color.rgb = BRIGHT_YELLOW
    icon.line.color.rgb = DEEP_RED
    icon.line.width = Pt(2)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(top), Inches(8.3), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    return title_box

def add_content_box(slide, content_lines, left, top, width, height, 
                    bg_color=CREAM, border_color=PRIMARY_RED, font_size=14):
    """Add styled content box with border"""
    # Background
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                  Inches(left), Inches(top), 
                                  Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = border_color
    box.line.width = Pt(2)
    
    # Content
    tb = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.15), 
                                   Inches(width - 0.3), Inches(height - 0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = BLACK
        p.space_after = Pt(6)

def add_numbered_bullet(slide, number, text, left, top, width=8):
    """Add numbered bullet with circle"""
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
                                     Inches(left), Inches(top), 
                                     Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY_RED
    circle.line.fill.background()
    
    # Number
    num_box = slide.shapes.add_textbox(Inches(left), Inches(top + 0.05), 
                                        Inches(0.4), Inches(0.35))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Text
    text_box = slide.shapes.add_textbox(Inches(left + 0.5), Inches(top), 
                                         Inches(width - 0.5), Inches(0.5))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = BLACK

def add_visual_divider(slide, top, left=0.5, width=9):
    """Add decorative divider line"""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                   Inches(left), Inches(top), 
                                   Inches(width), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = BRIGHT_YELLOW
    line.line.fill.background()

def create_infographic_slide(prs, title, items, subtitle=""):
    """Create slide with visual infographic layout"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_header(slide)
    add_decorative_footer(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = CREAM
    
    # Create visual boxes for items
    y_pos = 1.5
    x_positions = [0.5, 5.2]
    colors = [CREAM, LIGHT_YELLOW, RGBColor(255, 237, 213), RGBColor(254, 242, 242)]
    
    for i, (item_title, item_content) in enumerate(items):
        x = x_positions[i % 2]
        if i > 0 and i % 2 == 0:
            y_pos += 2.3
        
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(x), Inches(y_pos),
                                      Inches(4.3), Inches(2.1))
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.color.rgb = PRIMARY_RED if i % 2 == 0 else ORANGE
        box.line.width = Pt(2)
        
        # Number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(x + 0.1), Inches(y_pos + 0.1),
                                        Inches(0.35), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = PRIMARY_RED if i % 2 == 0 else ORANGE
        badge.line.fill.background()
        
        # Badge number
        badge_text = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y_pos + 0.13),
                                               Inches(0.35), Inches(0.3))
        tf = badge_text.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Item title
        item_title_box = slide.shapes.add_textbox(Inches(x + 0.55), Inches(y_pos + 0.15),
                                                   Inches(3.6), Inches(0.4))
        tf = item_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = item_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = DEEP_RED
        
        # Item content
        content_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y_pos + 0.6),
                                                Inches(3.9), Inches(1.4))
        tf = content_box.text_frame
        tf.word_wrap = True
        for line in item_content:
            p = tf.add_paragraph() if tf.paragraphs else tf.paragraphs[0]
            if len(tf.paragraphs) > 1 or not p.text:
                p = tf.add_paragraph() if p.text else p
            p.text = line
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(4)

def create_timeline_slide(prs, title, phases):
    """Create visual timeline slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_header(slide)
    add_decorative_footer(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Timeline line
    timeline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0.8), Inches(2.0),
                                       Inches(8.4), Inches(0.08))
    timeline.fill.solid()
    timeline.fill.fore_color.rgb = BRIGHT_YELLOW
    timeline.line.fill.background()
    
    # Phase boxes
    x_start = 1.0
    x_spacing = 1.7
    
    for i, (phase_name, phase_items) in enumerate(phases):
        x = x_start + (i * x_spacing)
        
        # Connector dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      Inches(x + 0.5), Inches(1.9),
                                      Inches(0.3), Inches(0.3))
        dot.fill.solid()
        dot.fill.fore_color.rgb = PRIMARY_RED
        dot.line.color.rgb = WHITE
        dot.line.width = Pt(2)
        
        # Week number
        week_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           Inches(x), Inches(2.5),
                                           Inches(1.3), Inches(0.5))
        week_box.fill.solid()
        week_box.fill.fore_color.rgb = PRIMARY_RED if i % 2 == 0 else ORANGE
        week_box.line.fill.background()
        
        week_text = slide.shapes.add_textbox(Inches(x), Inches(2.58),
                                              Inches(1.3), Inches(0.4))
        tf = week_text.text_frame
        p = tf.paragraphs[0]
        p.text = f"W{i}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Phase name
        name_box = slide.shapes.add_textbox(Inches(x - 0.1), Inches(3.15),
                                             Inches(1.5), Inches(0.6))
        tf = name_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = phase_name
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = DEEP_RED
        p.alignment = PP_ALIGN.CENTER
        
        # Content box
        content_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                              Inches(x - 0.1), Inches(3.8),
                                              Inches(1.5), Inches(2.5))
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = CREAM
        content_box.line.color.rgb = GOLD
        content_box.line.width = Pt(1.5)
        
        # Content items
        content_text = slide.shapes.add_textbox(Inches(x - 0.05), Inches(3.9),
                                                 Inches(1.4), Inches(2.3))
        tf = content_text.text_frame
        tf.word_wrap = True
        for item in phase_items:
            p = tf.add_paragraph() if tf.paragraphs else tf.paragraphs[0]
            if len(tf.paragraphs) > 1 or not p.text:
                p = tf.add_paragraph() if p.text else p
            p.text = "• " + item
            p.font.size = Pt(8)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(2)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ============ SLIDE 1 — CREATIVE COVER PAGE ============
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

# Full gradient background effect with shapes
bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(3.5))
bg1.fill.solid()
bg1.fill.fore_color.rgb = PRIMARY_RED
bg1.line.fill.background()

bg2 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.5), Inches(10), Inches(4))
bg2.fill.solid()
bg2.fill.fore_color.rgb = DEEP_RED
bg2.line.fill.background()

# Yellow wave accent
wave = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-1), Inches(3.2), Inches(12), Inches(0.6))
wave.fill.solid()
wave.fill.fore_color.rgb = BRIGHT_YELLOW
wave.line.fill.background()

# Decorative circles
circle1 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.5), Inches(1.2), Inches(1.2))
circle1.fill.solid()
circle1.fill.fore_color.rgb = BRIGHT_YELLOW
circle1.line.fill.background()

circle2 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.3), Inches(0.3), Inches(0.8), Inches(0.8))
circle2.fill.solid()
circle2.fill.fore_color.rgb = ORANGE
circle2.line.fill.background()

circle3 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.8), Inches(5.5), Inches(0.6), Inches(0.6))
circle3.fill.solid()
circle3.fill.fore_color.rgb = GOLD
circle3.line.fill.background()

# Main title
title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.2))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "XYZ Mobile App"
p.font.size = Pt(52)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Subtitle with background
sub_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                   Inches(1.5), Inches(3.0),
                                   Inches(7), Inches(0.8))
sub_box.fill.solid()
sub_box.fill.fore_color.rgb = WHITE
sub_box.line.color.rgb = GOLD
sub_box.line.width = Pt(3)

sub_text = slide1.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(7), Inches(0.6))
tf = sub_text.text_frame
p = tf.paragraphs[0]
p.text = "Performance & Release Management Diagnostic"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = DEEP_RED
p.alignment = PP_ALIGN.CENTER

# Company info
comp_box = slide1.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(8.4), Inches(0.6))
tf = comp_box.text_frame
p = tf.paragraphs[0]
p.text = "XYZ Company"
p.font.size = Pt(28)
p.font.color.rgb = BRIGHT_YELLOW
p.alignment = PP_ALIGN.CENTER

date_box = slide1.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(8.4), Inches(0.5))
tf = date_box.text_frame
p = tf.paragraphs[0]
p.text = "January 2026"
p.font.size = Pt(20)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 2 — Understanding Scope (Visual Cards) ============
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_header(slide2)
add_decorative_footer(slide2)

# Title with icon
add_title_with_icon(slide2, "Our Understanding of Scope", MSO_SHAPE.OVAL)

# Intro text
intro = slide2.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(8.8), Inches(0.4))
tf = intro.text_frame
p = tf.paragraphs[0]
p.text = "Client seeks a comprehensive diagnostic-driven assessment:"
p.font.size = Pt(14)
p.font.italic = True
p.font.color.rgb = DARK_GRAY

# Visual cards for each objective
cards = [
    ("Analyze", "Mobile app latency across Home, Insurance, Spend Track, Quiz & other key user flows"),
    ("Identify", "Root causes behind long load times (6 sec current vs 2-3 sec market benchmark)"),
    ("Understand", "App size inflation patterns (Android: 160MB→400MB+; iOS: 402MB)"),
    ("Determine", "Feasibility of moving to a monthly release cycle"),
    ("Recommend", "Fixes backed by measurable Root Cause Analysis - no assumptions"),
    ("Provide", "North Star performance vision to guide long-term optimization strategy")
]

y_pos = 1.9
for i, (head, desc) in enumerate(cards):
    x = 0.5 if i % 2 == 0 else 5.2
    if i > 0 and i % 2 == 0:
        y_pos += 1.8
    
    # Card background
    card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y_pos),
                                    Inches(4.3), Inches(1.6))
    card.fill.solid()
    card.fill.fore_color.rgb = CREAM if i % 2 == 0 else LIGHT_YELLOW
    card.line.color.rgb = PRIMARY_RED if i % 2 == 0 else ORANGE
    card.line.width = Pt(2)
    
    # Icon circle
    icon = slide2.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(x + 0.1), Inches(y_pos + 0.1),
                                    Inches(0.4), Inches(0.4))
    icon.fill.solid()
    icon.fill.fore_color.rgb = PRIMARY_RED if i % 2 == 0 else ORANGE
    icon.line.fill.background()
    
    # Heading
    head_box = slide2.shapes.add_textbox(Inches(x + 0.6), Inches(y_pos + 0.15),
                                          Inches(3.5), Inches(0.4))
    tf = head_box.text_frame
    p = tf.paragraphs[0]
    p.text = head
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DEEP_RED
    
    # Description
    desc_box = slide2.shapes.add_textbox(Inches(x + 0.15), Inches(y_pos + 0.6),
                                          Inches(4.0), Inches(0.9))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY

# ============ SLIDE 3 — Scope (Two-Column Visual) ============
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_header(slide3)
add_decorative_footer(slide3)

add_title_with_icon(slide3, "Scope of Diagnostic", MSO_SHAPE.OVAL)

# Left panel - AB Team
left_panel = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.4), Inches(1.4),
                                      Inches(4.6), Inches(5.3))
left_panel.fill.solid()
left_panel.fill.fore_color.rgb = RGBColor(254, 242, 242)
left_panel.line.color.rgb = PRIMARY_RED
left_panel.line.width = Pt(3)

# AB Team header
ab_header = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0.4), Inches(1.4),
                                     Inches(4.6), Inches(0.6))
ab_header.fill.solid()
ab_header.fill.fore_color.rgb = PRIMARY_RED
ab_header.line.fill.background()

ab_text = slide3.shapes.add_textbox(Inches(0.4), Inches(1.52), Inches(4.6), Inches(0.5))
tf = ab_text.text_frame
p = tf.paragraphs[0]
p.text = "🔧 AB Team Resources"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Team list
team_items = [
    ("Mobile Performance Lead", "Flutter Expert"),
    ("Mobile Performance Engineer", "Code Analysis"),
    ("Engineering Manager", "Program Management")
]
y = 2.2
for role, desc in team_items:
    # Role
    role_box = slide3.shapes.add_textbox(Inches(0.6), Inches(y), Inches(4.2), Inches(0.35))
    tf = role_box.text_frame
    p = tf.paragraphs[0]
    p.text = "▸ " + role
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DEEP_RED
    
    # Desc
    desc_box = slide3.shapes.add_textbox(Inches(0.9), Inches(y + 0.35), Inches(3.9), Inches(0.3))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_GRAY
    y += 0.8

# Activities header
act_header = slide3.shapes.add_textbox(Inches(0.6), Inches(4.5), Inches(4.2), Inches(0.4))
tf = act_header.text_frame
p = tf.paragraphs[0]
p.text = "📋 Key Activities"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = PRIMARY_RED

activities = ["Access setup: source code, UAT builds", "Journey & technical walkthrough",
              "Environment & build readiness confirmation"]
y = 4.9
for act in activities:
    act_box = slide3.shapes.add_textbox(Inches(0.7), Inches(y), Inches(4.0), Inches(0.35))
    tf = act_box.text_frame
    p = tf.paragraphs[0]
    p.text = "• " + act
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_GRAY
    y += 0.35

# Right panel - Client Inputs
right_panel = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(5.1), Inches(1.4),
                                       Inches(4.5), Inches(5.3))
right_panel.fill.solid()
right_panel.fill.fore_color.rgb = RGBColor(255, 251, 235)
right_panel.line.color.rgb = GOLD
right_panel.line.width = Pt(3)

# Client header
client_header = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         Inches(5.1), Inches(1.4),
                                         Inches(4.5), Inches(0.6))
client_header.fill.solid()
client_header.fill.fore_color.rgb = GOLD
client_header.line.fill.background()

client_text = slide3.shapes.add_textbox(Inches(5.1), Inches(1.52), Inches(4.5), Inches(0.5))
tf = client_text.text_frame
p = tf.paragraphs[0]
p.text = "📦 Inputs From Client"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Client inputs with icons
client_inputs = [
    ("📱", "Latest production build (APK/IPA)"),
    ("📊", "Access to Analytics, CMS"),
    ("📑", "API documentation"),
    ("🚀", "Release pipeline documentation"),
    ("🔌", "Third-party SDK list")
]
y = 2.2
for icon, item in client_inputs:
    item_box = slide3.shapes.add_textbox(Inches(5.3), Inches(y), Inches(4.1), Inches(0.5))
    tf = item_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon} {item}"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    y += 0.65

# ============ SLIDE 4 — North Star (Visual Benchmark) ============
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_header(slide4)
add_decorative_footer(slide4)

add_title_with_icon(slide4, "North Star Vision", MSO_SHAPE.OVAL)

# Subtitle
sub = slide4.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(8.8), Inches(0.4))
tf = sub.text_frame
p = tf.paragraphs[0]
p.text = "🎯 Industry Benchmarks & Target Goals"
p.font.size = Pt(14)
p.font.color.rgb = PRIMARY_RED
p.font.bold = True

# Benchmark cards
benchmarks = [
    ("⚡", "Screen Load Time", "2 sec", "Target", "6 sec", "Current", "60% improvement"),
    ("🔄", "Tab-Switch Latency", "150-250", "ms", "500+", "ms", "Visual feedback"),
    ("💾", "App Size Reduction", "30-40%", "Target", "Growing", "Current", "User retention"),
    ("🌐", "API Latency", "<150", "ms critical", "300+", "ms", "Core flows"),
    ("🎬", "Frame Stability", "<16", "ms/frame", "Janky", "Current", "Smooth UX"),
    ("📅", "Release Cadence", "Monthly", "Predictable", "Ad-hoc", "Current", "Consistent")
]

positions = [(0.4, 1.7), (3.5, 1.7), (6.6, 1.7), (0.4, 4.0), (3.5, 4.0), (6.6, 4.0)]

for i, (icon, label, target_val, target_unit, curr_val, curr_unit, benefit) in enumerate(benchmarks):
    x, y = positions[i]
    
    # Card
    card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y),
                                    Inches(2.8), Inches(2.15))
    card.fill.solid()
    card.fill.fore_color.rgb = CREAM if i % 2 == 0 else LIGHT_YELLOW
    card.line.color.rgb = PRIMARY_RED
    card.line.width = Pt(2)
    
    # Icon
    icon_box = slide4.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.1), Inches(0.5), Inches(0.5))
    tf = icon_box.text_frame
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(24)
    
    # Label
    label_box = slide4.shapes.add_textbox(Inches(x + 0.6), Inches(y + 0.15), Inches(2.0), Inches(0.4))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = DEEP_RED
    
    # Target value (highlighted)
    target_bg = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(x + 0.1), Inches(y + 0.6),
                                         Inches(1.2), Inches(0.6))
    target_bg.fill.solid()
    target_bg.fill.fore_color.rgb = RGBColor(220, 252, 231)
    target_bg.line.color.rgb = RGBColor(34, 197, 94)
    
    target_box = slide4.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.7), Inches(1.2), Inches(0.5))
    tf = target_box.text_frame
    p = tf.paragraphs[0]
    p.text = target_val
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(22, 101, 52)
    p.alignment = PP_ALIGN.CENTER
    
    target_unit_box = slide4.shapes.add_textbox(Inches(x + 1.35), Inches(y + 0.75), Inches(1.3), Inches(0.4))
    tf = target_unit_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{target_unit} (Target)"
    p.font.size = Pt(9)
    p.font.color.rgb = DARK_GRAY
    
    # Current value
    curr_box = slide4.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.25), Inches(2.5), Inches(0.3))
    tf = curr_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"vs {curr_val} {curr_unit}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(239, 68, 68)
    
    # Benefit
    benefit_box = slide4.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.6), Inches(2.5), Inches(0.4))
    tf = benefit_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"✓ {benefit}"
    p.font.size = Pt(9)
    p.font.italic = True
    p.font.color.rgb = DARK_GRAY

# Disclaimer
disc = slide4.shapes.add_textbox(Inches(0.4), Inches(6.3), Inches(9.2), Inches(0.4))
tf = disc.text_frame
p = tf.paragraphs[0]
p.text = "📌 Note: These are reference benchmarks only. Actual commitments determined post-RCA."
p.font.size = Pt(10)
p.font.italic = True
p.font.color.rgb = DARK_GRAY
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 5 — Assumptions (Visual) ============
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_header(slide5)
add_decorative_footer(slide5)

add_title_with_icon(slide5, "Key Assumptions", MSO_SHAPE.OVAL)

assumptions = [
    ("1", "Client Access", "All access (code, builds, dashboards) will be provided by the client in a timely manner"),
    ("2", "SDK Limitations", "Third-party SDK behavior and CMS limitations may restrict optimization scope"),
    ("3", "Backend Scope", "No changes to backend or CMS unless explicitly included in the engagement"),
    ("4", "RCA-Driven", "RCA outcomes will determine feasibility - not all items may be fixable"),
    ("5", "Data-Driven", "Recommendations will be measurable and derived from profiling & telemetry data"),
    ("6", "UX Stability", "Business-driven UI/UX changes out of scope unless mutually agreed"),
    ("7", "Advisory Role", "Release Management changes are advisory; client DevOps owns implementation")
]

y_pos = 1.5
for num, title, desc in assumptions:
    # Number circle
    circle = slide5.shapes.add_shape(MSO_SHAPE.OVAL,
                                      Inches(0.5), Inches(y_pos + 0.05),
                                      Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY_RED
    circle.line.fill.background()
    
    num_text = slide5.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.12), Inches(0.5), Inches(0.4))
    tf = num_text.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Title
    title_box = slide5.shapes.add_textbox(Inches(1.15), Inches(y_pos), Inches(3.0), Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = DEEP_RED
    
    # Description box
    desc_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(4.2), Inches(y_pos - 0.05),
                                        Inches(5.3), Inches(0.7))
    desc_box.fill.solid()
    desc_box.fill.fore_color.rgb = CREAM
    desc_box.line.color.rgb = GOLD
    
    desc_text = slide5.shapes.add_textbox(Inches(4.35), Inches(y_pos + 0.05), Inches(5.0), Inches(0.55))
    tf = desc_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_GRAY
    
    y_pos += 0.85

# ============ SLIDE 6 — Architecture (Visual Hub) ============
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_header(slide6)
add_decorative_footer(slide6)

add_title_with_icon(slide6, "Architecture & Design Focus", MSO_SHAPE.OVAL)

# Central hub visualization
# Center circle
center = slide6.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(4.0), Inches(3.0),
                                  Inches(2), Inches(2))
center.fill.solid()
center.fill.fore_color.rgb = PRIMARY_RED
center.line.color.rgb = BRIGHT_YELLOW
center.line.width = Pt(4)

center_text = slide6.shapes.add_textbox(Inches(4.0), Inches(3.7), Inches(2), Inches(0.8))
tf = center_text.text_frame
p = tf.paragraphs[0]
p.text = "Architecture\nAssessment"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Surrounding nodes
nodes = [
    ("Modular\nLayers", 1.5, 1.8, CREAM),
    ("API\nSequencing", 7.5, 1.8, LIGHT_YELLOW),
    ("Async\nRendering", 0.8, 4.0, RGBColor(254, 242, 242)),
    ("SDK\nFootprint", 8.2, 4.0, RGBColor(255, 251, 235)),
    ("Lazy\nLoading", 1.5, 5.8, RGBColor(236, 254, 255)),
    ("Asset\nCache", 7.5, 5.8, RGBColor(245, 243, 255))
]

for label, x, y, color in nodes:
    # Connection line
    line_x = 4.5 + (x - 4.5) * 0.5
    line_y = 3.5 + (y - 3.5) * 0.5
    
    # Node
    node = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y),
                                    Inches(1.8), Inches(1.0))
    node.fill.solid()
    node.fill.fore_color.rgb = color
    node.line.color.rgb = PRIMARY_RED
    node.line.width = Pt(2)
    
    # Label
    label_box = slide6.shapes.add_textbox(Inches(x), Inches(y + 0.25), Inches(1.8), Inches(0.6))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = DEEP_RED
    p.alignment = PP_ALIGN.CENTER

# Additional considerations at bottom
bottom_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(2.5), Inches(6.4),
                                      Inches(5), Inches(0.6))
bottom_box.fill.solid()
bottom_box.fill.fore_color.rgb = CREAM
bottom_box.line.color.rgb = GOLD

bottom_text = slide6.shapes.add_textbox(Inches(2.6), Inches(6.55), Inches(4.8), Inches(0.4))
tf = bottom_text.text_frame
p = tf.paragraphs[0]
p.text = "🎯 Plus: Separation of concerns, Release governance, Branching strategy review"
p.font.size = Pt(11)
p.font.color.rgb = DARK_GRAY
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 7 — Diagnostic View (Layered) ============
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_header(slide7)
add_decorative_footer(slide7)

add_title_with_icon(slide7, "Proposed Diagnostic Architecture", MSO_SHAPE.FLOWCHART_PROCESS)

# Layered architecture visualization
layers = [
    ("📱 App Frontend", "Flutter architecture deep-dive", RGBColor(254, 242, 242), 0.5),
    ("🌐 API Gateway", "Interaction patterns & latency", RGBColor(255, 251, 235), 1.4),
    ("📦 CMS Modules", "Content-driven module analysis", RGBColor(236, 254, 255), 2.3),
    ("🔌 SDK Integrations", "Fly, MarTech, Payments, Firebase", RGBColor(245, 243, 255), 3.2),
    ("📊 Telemetry", "Performance data flows", RGBColor(255, 241, 242), 4.1),
    ("🚀 CI/CD Pipeline", "Release workflows & automation", RGBColor(255, 247, 237), 5.0)
]

for icon_label, desc, color, y_offset in layers:
    # Layer bar
    layer = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(2), Inches(1.5 + y_offset),
                                     Inches(6), Inches(0.75))
    layer.fill.solid()
    layer.fill.fore_color.rgb = color
    layer.line.color.rgb = PRIMARY_RED
    layer.line.width = Pt(2)
    
    # Icon + Label
    label_box = slide7.shapes.add_textbox(Inches(2.2), Inches(1.58 + y_offset),
                                           Inches(3.5), Inches(0.6))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = icon_label
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DEEP_RED
    
    # Description
    desc_box = slide7.shapes.add_textbox(Inches(5.5), Inches(1.65 + y_offset),
                                          Inches(2.3), Inches(0.55))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_GRAY
    
    # Arrow down (except last)
    if y_offset < 5.0:
        arrow = slide7.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                         Inches(4.8), Inches(2.28 + y_offset),
                                         Inches(0.4), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD
        arrow.line.fill.background()

# ============ SLIDE 8 — Timeline (Visual Flow) ============
phases = [
    ("Setup", ["Access provision", "Environment", "Build validation"]),
    ("Profiling", ["Load benchmarks", "API mapping", "Size breakdown", "Cache profiling"]),
    ("RCA", ["Root-cause ID", "Fixable analysis", "SDK constraints", "Pipeline review"]),
    ("Validation", ["Joint walkthrough", "Team validation", "Feasibility check", "Release model"]),
    ("Report", ["Target alignment", "Uplift range", "Fixes list", "Execution plan"])
]
create_timeline_slide(prs, "Diagnostic Approach: Week 0-4", phases)

# ============ SLIDE 9 — Team (Visual Org) ============
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_header(slide9)
add_decorative_footer(slide9)

add_title_with_icon(slide9, "Expert Team Structure", MSO_SHAPE.OVAL)

# Team visualization with photos placeholder
team_data = [
    ("👨‍💻", "Mobile Performance Lead", "Flutter Expert", "Profiling, rendering, optimization",
     "1", PRIMARY_RED),
    ("👩‍💻", "Mobile Engineers", "Code Analysts", "Code analysis, architecture assessment",
     "2", ORANGE),
    ("🎯", "Solution Architect", "Strategic Lead", "Engineering Mgmt, Program Mgmt, Architecture",
     "1", GOLD)
]

x_positions = [0.5, 3.7, 6.9]
for i, (icon, role, subtitle, resp, count, color) in enumerate(team_data):
    x = x_positions[i]
    
    # Card
    card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(1.8),
                                    Inches(2.6), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = CREAM
    card.line.color.rgb = color
    card.line.width = Pt(3)
    
    # Avatar circle
    avatar = slide9.shapes.add_shape(MSO_SHAPE.OVAL,
                                      Inches(x + 0.8), Inches(2.0),
                                      Inches(1), Inches(1))
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = color
    avatar.line.color.rgb = WHITE
    avatar.line.width = Pt(3)
    
    # Icon
    icon_box = slide9.shapes.add_textbox(Inches(x + 0.8), Inches(2.2), Inches(1), Inches(0.6))
    tf = icon_box.text_frame
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(36)
    p.alignment = PP_ALIGN.CENTER
    
    # Role
    role_box = slide9.shapes.add_textbox(Inches(x + 0.1), Inches(3.2), Inches(2.4), Inches(0.5))
    tf = role_box.text_frame
    p = tf.paragraphs[0]
    p.text = role
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = DEEP_RED
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide9.shapes.add_textbox(Inches(x + 0.1), Inches(3.7), Inches(2.4), Inches(0.4))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    
    # Count badge
    badge = slide9.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(x + 1.9), Inches(1.6),
                                     Inches(0.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    
    count_text = slide9.shapes.add_textbox(Inches(x + 1.9), Inches(1.68), Inches(0.5), Inches(0.4))
    tf = count_text.text_frame
    p = tf.paragraphs[0]
    p.text = count
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Responsibility
    resp_bg = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x + 0.2), Inches(4.2),
                                       Inches(2.2), Inches(1.8))
    resp_bg.fill.solid()
    resp_bg.fill.fore_color.rgb = WHITE
    resp_bg.line.color.rgb = GOLD
    
    resp_text = slide9.shapes.add_textbox(Inches(x + 0.3), Inches(4.35), Inches(2.0), Inches(1.6))
    tf = resp_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = resp
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

# Total team
total_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(3), Inches(6.5),
                                     Inches(4), Inches(0.5))
total_box.fill.solid()
total_box.fill.fore_color.rgb = PRIMARY_RED
total_box.line.fill.background()

total_text = slide9.shapes.add_textbox(Inches(3), Inches(6.6), Inches(4), Inches(0.4))
tf = total_text.text_frame
p = tf.paragraphs[0]
p.text = "🚀 Total Team: 3 Members (Diagnostic Phase)"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 10 — Gantt (Visual Bars) ============
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_header(slide10)
add_decorative_footer(slide10)

add_title_with_icon(slide10, "Project Timeline (Gantt View)", MSO_SHAPE.RECTANGLE)

# Gantt chart
weeks = ["W0", "W1", "W2", "W3", "W4"]
tasks = [
    ("Setup & Access", 0, 1, PRIMARY_RED),
    ("Profiling & Analysis", 1, 1, ORANGE),
    ("RCA & Observations", 2, 1, GOLD),
    ("Discussions & Validation", 3, 1, RGBColor(34, 197, 94)),
    ("Final Report", 4, 1, RGBColor(59, 130, 246))
]

# Week headers
for i, week in enumerate(weeks):
    header = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(4 + i * 1.15), Inches(1.5),
                                       Inches(1.05), Inches(0.5))
    header.fill.solid()
    header.fill.fore_color.rgb = DEEP_RED
    header.line.fill.background()
    
    week_text = slide10.shapes.add_textbox(Inches(4 + i * 1.15), Inches(1.58),
                                            Inches(1.05), Inches(0.4))
    tf = week_text.text_frame
    p = tf.paragraphs[0]
    p.text = week
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# Task bars
y_pos = 2.3
for task_name, start, duration, color in tasks:
    # Task label
    label = slide10.shapes.add_textbox(Inches(0.3), Inches(y_pos), Inches(3.5), Inches(0.5))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = task_name
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY
    
    # Bar
    bar = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(4 + start * 1.15), Inches(y_pos + 0.05),
                                    Inches(duration * 1.05), Inches(0.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.color.rgb = WHITE
    bar.line.width = Pt(2)
    
    y_pos += 0.7

# Milestones
milestones = [
    (4.55, "Kickoff"),
    (5.7, "Data Ready"),
    (6.85, "Findings"),
    (8.0, "Validation"),
    (9.15, "Delivery")
]

for x, label in milestones:
    # Diamond marker
    diamond = slide10.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                        Inches(x - 0.1), Inches(6.0),
                                        Inches(0.2), Inches(0.2))
    diamond.fill.solid()
    diamond.fill.fore_color.rgb = GOLD
    diamond.line.fill.background()
    
    # Label
    m_label = slide10.shapes.add_textbox(Inches(x - 0.4), Inches(6.3), Inches(0.8), Inches(0.4))
    tf = m_label.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(8)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 11 — Commercial (Visual Pricing) ============
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_header(slide11)
add_decorative_footer(slide11)

add_title_with_icon(slide11, "Investment Structure", MSO_SHAPE.OVAL)

# Phase boxes - side by side
# Diagnostic Phase
diag_box = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.5), Inches(1.6),
                                     Inches(4.4), Inches(4.8))
diag_box.fill.solid()
diag_box.fill.fore_color.rgb = RGBColor(254, 242, 242)
diag_box.line.color.rgb = PRIMARY_RED
diag_box.line.width = Pt(3)

# Header
diag_header = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.5), Inches(1.6),
                                        Inches(4.4), Inches(0.8))
diag_header.fill.solid()
diag_header.fill.fore_color.rgb = PRIMARY_RED
diag_header.line.fill.background()

diag_title = slide11.shapes.add_textbox(Inches(0.5), Inches(1.75), Inches(4.4), Inches(0.6))
tf = diag_title.text_frame
p = tf.paragraphs[0]
p.text = "🔍 DIAGNOSTIC PHASE"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Duration badge
duration_badge = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           Inches(1.8), Inches(2.55),
                                           Inches(1.8), Inches(0.5))
duration_badge.fill.solid()
duration_badge.fill.fore_color.rgb = BRIGHT_YELLOW
duration_badge.line.fill.background()

dur_text = slide11.shapes.add_textbox(Inches(1.8), Inches(2.62), Inches(1.8), Inches(0.4))
tf = dur_text.text_frame
p = tf.paragraphs[0]
p.text = "⏱ 4 WEEKS"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = DEEP_RED
p.alignment = PP_ALIGN.CENTER

# Features
diag_features = [
    "✓ Fixed fee structure",
    "✓ 4-5 resources (1 month FTE)",
    "✓ Comprehensive analysis",
    "✓ Root cause documentation",
    "✓ Release advisory report"
]
y = 3.2
for feat in diag_features:
    feat_box = slide11.shapes.add_textbox(Inches(0.8), Inches(y), Inches(3.8), Inches(0.4))
    tf = feat_box.text_frame
    p = tf.paragraphs[0]
    p.text = feat
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY
    y += 0.5

# Execution Phase
exec_box = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(5.1), Inches(1.6),
                                     Inches(4.4), Inches(4.8))
exec_box.fill.solid()
exec_box.fill.fore_color.rgb = RGBColor(255, 251, 235)
exec_box.line.color.rgb = GOLD
exec_box.line.width = Pt(3)

exec_header = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(5.1), Inches(1.6),
                                        Inches(4.4), Inches(0.8))
exec_header.fill.solid()
exec_header.fill.fore_color.rgb = GOLD
exec_header.line.fill.background()

exec_title = slide11.shapes.add_textbox(Inches(5.1), Inches(1.75), Inches(4.4), Inches(0.6))
tf = exec_title.text_frame
p = tf.paragraphs[0]
p.text = "🚀 EXECUTION PHASE"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

exec_badge = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(6.4), Inches(2.55),
                                       Inches(1.8), Inches(0.5))
exec_badge.fill.solid()
exec_badge.fill.fore_color.rgb = ORANGE
exec_badge.line.fill.background()

exec_badge_text = slide11.shapes.add_textbox(Inches(6.4), Inches(2.62), Inches(1.8), Inches(0.4))
tf = exec_badge_text.text_frame
p = tf.paragraphs[0]
p.text = "📋 TBD"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

exec_features = [
    "⏳ Estimation post-diagnostic",
    "⏳ Dependent on fixable scope",
    "⏳ Detailed sprint planning",
    "⏳ Resource allocation",
    "⏳ Timeline confirmation"
]
y = 3.2
for feat in exec_features:
    feat_box = slide11.shapes.add_textbox(Inches(5.4), Inches(y), Inches(3.8), Inches(0.4))
    tf = feat_box.text_frame
    p = tf.paragraphs[0]
    p.text = feat
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY
    y += 0.5

# Arrow between
arrow = slide11.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                  Inches(4.7), Inches(3.5),
                                  Inches(0.6), Inches(0.4))
arrow.fill.solid()
arrow.fill.fore_color.rgb = PRIMARY_RED
arrow.line.fill.background()

# ============ SLIDE 12 — Risks (Visual Warning) ============
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_header(slide12)
add_decorative_footer(slide12)

add_title_with_icon(slide12, "Risk Factors & Dependencies", MSO_SHAPE.OVAL)

risks = [
    ("⚠️", "Third-Party SDKs", "External SDK behavior may limit optimization options"),
    ("📡", "CMS Constraints", "Content Management System payload restrictions"),
    ("⏰", "API Dependencies", "Launch-time API calls blocking user experience"),
    ("📱", "Device Fragmentation", "Low-RAM device behavior variations"),
    ("🚀", "Release Maturity", "Current release process capabilities"),
    ("🔧", "Environment Access", "UAT/Production environment availability")
]

positions = [(0.5, 1.6), (5.2, 1.6), (0.5, 3.6), (5.2, 3.6), (0.5, 5.6), (5.2, 5.6)]

for i, (icon, title, desc) in enumerate(risks):
    x, y = positions[i]
    
    # Risk card
    card = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(x), Inches(y),
                                     Inches(4.3), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(254, 242, 242)
    card.line.color.rgb = PRIMARY_RED
    card.line.width = Pt(2)
    
    # Warning stripe
    stripe = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(x), Inches(y),
                                       Inches(0.15), Inches(1.8))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ORANGE
    stripe.line.fill.background()
    
    # Icon
    icon_box = slide12.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.1), Inches(0.6), Inches(0.6))
    tf = icon_box.text_frame
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(24)
    
    # Title
    title_box = slide12.shapes.add_textbox(Inches(x + 0.9), Inches(y + 0.15), Inches(3.2), Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = DEEP_RED
    
    # Description
    desc_box = slide12.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.75), Inches(3.9), Inches(0.9))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_GRAY

# ============ SLIDE 13 — Final Outcome (Celebration) ============
slide13 = prs.slides.add_slide(prs.slide_layouts[6])

# Gradient background
bg_top = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(4))
bg_top.fill.solid()
bg_top.fill.fore_color.rgb = PRIMARY_RED
bg_top.line.fill.background()

bg_bottom = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4), Inches(10), Inches(3.5))
bg_bottom.fill.solid()
bg_bottom.fill.fore_color.rgb = CREAM
bg_bottom.line.fill.background()

# Yellow wave
wave2 = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-1), Inches(3.6), Inches(12), Inches(0.5))
wave2.fill.solid()
wave2.fill.fore_color.rgb = BRIGHT_YELLOW
wave2.line.fill.background()

add_decorative_footer(slide13)

# Title
final_title = slide13.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
tf = final_title.text_frame
p = tf.paragraphs[0]
p.text = "🎯 Final Deliverables"
p.font.size = Pt(38)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Subtitle
final_sub = slide13.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
tf = final_sub.text_frame
p = tf.paragraphs[0]
p.text = "Your comprehensive, data-backed performance roadmap"
p.font.size = Pt(16)
p.font.color.rgb = CREAM
p.alignment = PP_ALIGN.CENTER

# Deliverable cards
deliverables = [
    ("✅", "Improvement Plan", "What CAN be improved with estimated impact"),
    ("❌", "Constraints Doc", "What CANNOT be improved & why (SDK/CMS limits)"),
    ("📈", "Uplift Range", "Expected performance gains per optimization area"),
    ("👥", "Execution Plan", "Team structure & timeline for implementation"),
    ("📅", "Release Readiness", "Monthly release train feasibility assessment")
]

x_positions = [0.4, 2.1, 3.8, 5.5, 7.2]
for i, (icon, title, desc) in enumerate(deliverables):
    x = x_positions[i]
    
    # Card
    card = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(x), Inches(2.2),
                                     Inches(1.5), Inches(2.3))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = GOLD
    card.line.width = Pt(2)
    
    # Shadow effect
    shadow = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x + 0.05), Inches(2.25),
                                       Inches(1.5), Inches(2.3))
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(0, 0, 0)
    shadow.line.fill.background()
    # Move shadow behind (would need to reorder, visual only)
    
    # Icon
    icon_box = slide13.shapes.add_textbox(Inches(x + 0.5), Inches(2.4), Inches(0.5), Inches(0.5))
    tf = icon_box.text_frame
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(28)
    
    # Title
    title_box = slide13.shapes.add_textbox(Inches(x + 0.1), Inches(2.95), Inches(1.3), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = DEEP_RED
    p.alignment = PP_ALIGN.CENTER
    
    # Desc
    desc_box = slide13.shapes.add_textbox(Inches(x + 0.1), Inches(3.6), Inches(1.3), Inches(0.8))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(8)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

# Closing statement
closing = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1), Inches(5.0),
                                    Inches(8), Inches(1.0))
closing.fill.solid()
closing.fill.fore_color.rgb = PRIMARY_RED
closing.line.fill.background()

closing_text = slide13.shapes.add_textbox(Inches(1), Inches(5.3), Inches(8), Inches(0.6))
tf = closing_text.text_frame
p = tf.paragraphs[0]
p.text = "🚀 North Star Aligned • Data-Backed • Feasible & Actionable"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Save with timestamp
timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
output_name = f'XYZ_Mobile_App_Diagnostic_Creative_{timestamp}.pptx'
output_path = f'/Users/Mudassar.Hakim/tempfiles/{output_name}'

prs.save(output_path)
print(f"🎨 Creative presentation created successfully: {output_name}")
print(f"📊 Features: Infographics, visual timelines, team cards, benchmark grids")
