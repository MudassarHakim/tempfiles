from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.dml.color import RGBColor

# Colors - Red and Yellow theme
RED = RGBColor(200, 30, 30)
DARK_RED = RGBColor(139, 0, 0)
YELLOW = RGBColor(255, 200, 0)
DARK_YELLOW = RGBColor(218, 165, 32)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

def add_title_shape(slide, text, top=0.3, font_size=32, color=WHITE):
    """Add a title text box"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    return title_box

def add_body_text(slide, text_lines, top=1.2, left=0.5, width=9, font_size=14):
    """Add body text with bullet points"""
    body_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = body_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = BLACK
        p.level = 0
        p.space_after = Pt(8)
    return body_box

def add_header_bar(slide, color=RED):
    """Add a colored header bar at top"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_accent_bar(slide, top=1.1, color=YELLOW):
    """Add accent bar below header"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(top), Inches(10), Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_footer_bar(slide, color=DARK_RED):
    """Add footer bar"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.2), Inches(10), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def create_content_slide(prs, title, content_lines):
    """Create a standard content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_header_bar(slide)
    add_accent_bar(slide)
    add_footer_bar(slide)
    add_title_shape(slide, title, top=0.15, font_size=28)
    add_body_text(slide, content_lines, top=1.4)
    return slide

def create_two_column_slide(prs, title, left_title, left_content, right_title, right_content):
    """Create a two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide)
    add_accent_bar(slide)
    add_footer_bar(slide)
    add_title_shape(slide, title, top=0.15, font_size=28)
    
    # Left column title
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.3), Inches(0.4))
    tf = left_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_RED
    
    # Left column content
    add_body_text(slide, left_content, top=1.7, left=0.5, width=4.3, font_size=12)
    
    # Right column title
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(0.4))
    tf = right_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_RED
    
    # Right column content
    right_body = slide.shapes.add_textbox(Inches(5.2), Inches(1.7), Inches(4.3), Inches(5))
    tf = right_body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK
        p.space_after = Pt(6)
    
    return slide

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ============ SLIDE 1 — Cover Page ============
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
# Full red background
bg_shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = RED
bg_shape.line.fill.background()

# Yellow accent stripe
accent = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3), Inches(10), Inches(0.15))
accent.fill.solid()
accent.fill.fore_color.rgb = YELLOW
accent.line.fill.background()

# Title
title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "XYZ Mobile App"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Subtitle
sub_box = slide1.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(0.8))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "Performance & Release Management Diagnostic"
p.font.size = Pt(28)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Company
comp_box = slide1.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.5))
tf = comp_box.text_frame
p = tf.paragraphs[0]
p.text = "XYZ Company"
p.font.size = Pt(20)
p.font.color.rgb = YELLOW
p.alignment = PP_ALIGN.CENTER

# Date
date_box = slide1.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(9), Inches(0.5))
tf = date_box.text_frame
p = tf.paragraphs[0]
p.text = "January 2026"
p.font.size = Pt(18)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 2 — Our Understanding of Scope ============
slide2_content = [
    "Client seeks a diagnostic-driven assessment to:",
    "",
    "• Analyze mobile app latency across Home, Insurance, Spend Track, Quiz & other flows",
    "• Identify root causes behind long load times (6 seconds vs market 2–3 sec benchmark)",
    "• Understand app size inflation (Android: 160MB → 400+MB installed; iOS: 402MB)",
    "• Determine feasibility of moving to a monthly release cycle",
    "• Recommend fixes backed by measurable RCA (no assumptions)",
    "• Provide a North Star performance vision to guide long-term optimization"
]
create_content_slide(prs, "Our Understanding of Scope", slide2_content)

# ============ SLIDE 3 — Scope of Diagnostic ============
left_content = [
    "AB Team:",
    "• Mobile Performance Lead (Flutter)",
    "• Mobile Performance Engineer",
    "• Engineering Manager",
    "",
    "Activities:",
    "• Access setup: source code, UAT builds",
    "• Journey & technical walkthrough",
    "• Environment & build readiness confirmation"
]
right_content = [
    "Inputs Needed From Client:",
    "• Latest production build (APK/IPA)",
    "• Access to Analytics, CMS",
    "• API documentation",
    "• Release pipeline documentation",
    "• Third-party SDK list"
]
create_two_column_slide(prs, "Scope of Diagnostic (Mapped to Reference Structure)", 
                        "Project Start – Pre-Requisite", left_content,
                        "Diagnostic Pre-Requisite", right_content)

# ============ SLIDE 4 — North Star Vision ============
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide4)
add_accent_bar(slide4)
add_footer_bar(slide4)
add_title_shape(slide4, "North Star Vision (Benchmarking)", top=0.15, font_size=28)

# Subtitle
sub = slide4.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.4))
tf = sub.text_frame
p = tf.paragraphs[0]
p.text = "Performance North Star (Industry Benchmarks)"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_RED

benchmarks = [
    "• Primary screen load time: 2 seconds (Market), Current 6 seconds",
    "• Tab-switch latency: 150–250 ms",
    "• App Size Target: 30%–40% reduction",
    "• API latency goal: <150 ms for critical flows",
    "• Rendering frame stability: <16ms per frame",
    "• Release cadence: Predictable Monthly Release Train",
    "",
    "Note: These are reference benchmarks only, not commitments until RCA is completed."
]
add_body_text(slide4, benchmarks, top=1.8)

# ============ SLIDE 5 — Assumptions ============
slide5_content = [
    "• All access (code, builds, dashboards) will be provided by the client",
    "• Third-party SDK behavior and CMS limitations may restrict optimization",
    "• No changes to backend or CMS unless explicitly included",
    "• RCA outcomes will determine feasibility of performance enhancements",
    "• Recommendations will be measurable and derived from profiling & data",
    "• Any business-driven UI/UX changes are out of scope unless mutually agreed",
    "• Release Management changes are advisory; implementation may require client DevOps involvement"
]
create_content_slide(prs, "Assumptions", slide5_content)

# ============ SLIDE 6 — Architecture & Design Considerations ============
slide6_content = [
    "• Modular, layered architecture assessment",
    "• API sequencing, dependency mapping",
    "• Asynchronous vs synchronous rendering optimization",
    "• Third-party SDK footprint & load behavior",
    "• Lazy-loading feasibility",
    "• Asset compression & caching strategies",
    "• Separation of concerns for future scalability",
    "• Release governance & branching strategy review"
]
create_content_slide(prs, "Architecture & Design Considerations", slide6_content)

# ============ SLIDE 7 — Proposed Diagnostic Architecture View ============
slide7_content = [
    "Includes review of:",
    "",
    "• App frontend architecture (Flutter)",
    "• API Gateway interactions",
    "• CMS-driven modules",
    "• Third-party SDK integrations (Fly, MarTech, Payments, Firebase)",
    "• Performance telemetry flows",
    "• Release pipeline & CI/CD workflows"
]
create_content_slide(prs, "Proposed Diagnostic Architecture View", slide7_content)

# ============ SLIDE 8 — Our Approach (Week 0 to Week 4) ============
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide8)
add_accent_bar(slide8)
add_footer_bar(slide8)
add_title_shape(slide8, "Our Approach (Diagnostic) – Week 0 to Week 4", top=0.15, font_size=26)

# Week boxes - arranged in grid
weeks_data = [
    ("Week 0 – Setup & Access", ["Access provisioning", "Environment setup", "Build validation"]),
    ("Week 1 – Profiling & Analysis", ["Load-time benchmarking", "API call mapping & latency analysis", "App size/composition breakdown", "Cache & render-path profiling"]),
    ("Week 2 – RCA & Observations", ["Root-cause identification (H/M/L)", "Fixable vs dependency-driven issues", "Third-party SDK constraints", "Release pipeline maturity assessment"]),
    ("Week 3 – Discussions", ["Joint walkthrough of findings", "Validation with client teams", "Feasibility confirmation", "Monthly release model alignment"]),
    ("Week 4 – Final Report", ["North Star vs Achievable Targets", "Performance uplift range", "Recommended fixes", "Execution Plan & Sizing"])
]

positions = [(0.3, 1.4), (3.5, 1.4), (6.7, 1.4), (1.9, 4.2), (5.1, 4.2)]

for i, (week_title, items) in enumerate(weeks_data):
    x, y = positions[i]
    # Box background
    box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3), Inches(2.4))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 245, 200)  # Light yellow
    box.line.color.rgb = DARK_RED
    box.line.width = Pt(2)
    
    # Week title
    title_tb = slide8.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.1), Inches(2.8), Inches(0.4))
    tf = title_tb.text_frame
    p = tf.paragraphs[0]
    p.text = week_title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = DARK_RED
    
    # Items
    items_tb = slide8.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.5), Inches(2.8), Inches(1.8))
    tf = items_tb.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(9)
        p.font.color.rgb = BLACK

# ============ SLIDE 9 — Team Structure ============
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide9)
add_accent_bar(slide9)
add_footer_bar(slide9)
add_title_shape(slide9, "Team Structure", top=0.15, font_size=28)

# Table header
header_y = 1.5
col_widths = [3.5, 0.8, 4.5]
col_x = [0.6, 4.1, 5.0]
headers = ["Role", "Count", "Responsibility"]

for i, (header, x, w) in enumerate(zip(headers, col_x, col_widths)):
    box = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(header_y), Inches(w), Inches(0.5))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_RED
    box.line.fill.background()
    
    tb = slide9.shapes.add_textbox(Inches(x), Inches(header_y + 0.1), Inches(w), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = header
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# Table rows
rows_data = [
    ("Mobile Performance Lead", "1", "Profiling, rendering, app load optimization"),
    ("Mobile Engineer", "2", "Code analysis, architecture assessment"),
    ("Solution Architect", "1", "Engineering Manager, Program Management, Architecture")
]

row_y = 2.0
for role, count, resp in rows_data:
    row_y += 0.6
    # Role cell
    box1 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(col_x[0]), Inches(row_y), Inches(col_widths[0]), Inches(0.55))
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(255, 245, 200)
    box1.line.color.rgb = DARK_RED
    
    tb1 = slide9.shapes.add_textbox(Inches(col_x[0] + 0.1), Inches(row_y + 0.1), Inches(col_widths[0] - 0.2), Inches(0.4))
    tf = tb1.text_frame
    p = tf.paragraphs[0]
    p.text = role
    p.font.size = Pt(12)
    p.font.color.rgb = BLACK
    
    # Count cell
    box2 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(col_x[1]), Inches(row_y), Inches(col_widths[1]), Inches(0.55))
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(255, 245, 200)
    box2.line.color.rgb = DARK_RED
    
    tb2 = slide9.shapes.add_textbox(Inches(col_x[1]), Inches(row_y + 0.1), Inches(col_widths[1]), Inches(0.4))
    tf = tb2.text_frame
    p = tf.paragraphs[0]
    p.text = count
    p.font.size = Pt(12)
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER
    
    # Responsibility cell
    box3 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(col_x[2]), Inches(row_y), Inches(col_widths[2]), Inches(0.55))
    box3.fill.solid()
    box3.fill.fore_color.rgb = RGBColor(255, 245, 200)
    box3.line.color.rgb = DARK_RED
    
    tb3 = slide9.shapes.add_textbox(Inches(col_x[2] + 0.1), Inches(row_y + 0.1), Inches(col_widths[2] - 0.2), Inches(0.4))
    tf = tb3.text_frame
    p = tf.paragraphs[0]
    p.text = resp
    p.font.size = Pt(11)
    p.font.color.rgb = BLACK

# Total row
row_y += 0.8
total_box = slide9.shapes.add_textbox(Inches(0.6), Inches(row_y), Inches(8.8), Inches(0.5))
tf = total_box.text_frame
p = tf.paragraphs[0]
p.text = "Total Team: 3 Members (Diagnostic Phase)"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK_RED

# ============ SLIDE 10 — Gantt Timeline ============
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide10)
add_accent_bar(slide10)
add_footer_bar(slide10)
add_title_shape(slide10, "Gantt Timeline (Week-Level)", top=0.15, font_size=28)

# Gantt chart visualization
gantt_data = [
    ("Week 0 – Setup & Access", 0, 1),
    ("Week 1 – Profiling & Analysis", 1, 1),
    ("Week 2 – RCA & Observations", 2, 1),
    ("Week 3 – Discussions & Verifications", 3, 1),
    ("Week 4 – Final Report", 4, 1)
]

# Week headers
for i in range(5):
    header_box = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4 + i * 1.1), Inches(1.5), Inches(1.05), Inches(0.4))
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = DARK_RED
    header_box.line.fill.background()
    
    htb = slide10.shapes.add_textbox(Inches(4 + i * 1.1), Inches(1.55), Inches(1.05), Inches(0.35))
    tf = htb.text_frame
    p = tf.paragraphs[0]
    p.text = f"W{i}"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# Gantt bars
bar_y = 2.0
for task, start, duration in gantt_data:
    bar_y += 0.7
    # Task label
    label_box = slide10.shapes.add_textbox(Inches(0.3), Inches(bar_y), Inches(3.5), Inches(0.5))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = task
    p.font.size = Pt(11)
    p.font.color.rgb = BLACK
    
    # Gantt bar
    bar = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4 + start * 1.1), Inches(bar_y + 0.05), Inches(duration * 1.05), Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = YELLOW
    bar.line.color.rgb = DARK_RED
    bar.line.width = Pt(1.5)

# ============ SLIDE 11 — Commercial Structure ============
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide11)
add_accent_bar(slide11)
add_footer_bar(slide11)
add_title_shape(slide11, "Commercial Structure", top=0.15, font_size=28)

# Diagnostic Phase box
diag_box = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(4.3), Inches(2.5))
diag_box.fill.solid()
diag_box.fill.fore_color.rgb = RGBColor(255, 245, 200)
diag_box.line.color.rgb = DARK_RED
diag_box.line.width = Pt(2)

diag_title = slide11.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(3.9), Inches(0.5))
tf = diag_title.text_frame
p = tf.paragraphs[0]
p.text = "Diagnostic Phase (4 Weeks)"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = DARK_RED

diag_content = slide11.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(3.9), Inches(1.5))
tf = diag_content.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "• Fixed fee (based on 4–5 resources for 1 month equivalent)"
p.font.size = Pt(12)
p.font.color.rgb = BLACK
p = tf.add_paragraph()
p.text = "• Covers analysis, RCA, reporting, release advisory"
p.font.size = Pt(12)
p.font.color.rgb = BLACK

# Execution Phase box
exec_box = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.5), Inches(4.3), Inches(2.5))
exec_box.fill.solid()
exec_box.fill.fore_color.rgb = RGBColor(255, 245, 200)
exec_box.line.color.rgb = DARK_RED
exec_box.line.width = Pt(2)

exec_title = slide11.shapes.add_textbox(Inches(5.4), Inches(1.6), Inches(3.9), Inches(0.5))
tf = exec_title.text_frame
p = tf.paragraphs[0]
p.text = "Execution Phase"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = DARK_RED

exec_content = slide11.shapes.add_textbox(Inches(5.4), Inches(2.2), Inches(3.9), Inches(1.5))
tf = exec_content.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "• To be estimated based on diagnostic output"
p.font.size = Pt(12)
p.font.color.rgb = BLACK
p = tf.add_paragraph()
p.text = "• Dependent on size of fixable items"
p.font.size = Pt(12)
p.font.color.rgb = BLACK

# ============ SLIDE 12 — Risks & Dependencies ============
slide12_content = [
    "• Third-party SDK limitations",
    "• CMS payload constraints",
    "• Launch-time API dependencies",
    "• Device fragmentation & low-RAM behavior",
    "• Release process maturity",
    "• Environment availability"
]
create_content_slide(prs, "Risks & Dependencies", slide12_content)

# ============ SLIDE 13 — Final Outcome ============
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide13)
add_accent_bar(slide13)
add_footer_bar(slide13)
add_title_shape(slide13, "Final Outcome", top=0.15, font_size=28)

intro = slide13.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(0.5))
tf = intro.text_frame
p = tf.paragraphs[0]
p.text = "A North-Star aligned, data-backed, feasible performance roadmap including:"
p.font.size = Pt(16)
p.font.color.rgb = DARK_RED
p.font.bold = True

outcomes = [
    "• What can be improved",
    "• What cannot be improved",
    "• Expected uplift range",
    "• Team & timeline for execution",
    "• Monthly release readiness assessment"
]
add_body_text(slide13, outcomes, top=2.0, font_size=16)

# Save the presentation
prs.save('/Users/Mudassar.Hakim/tempfiles/XYZ_Mobile_App_Diagnostic.pptx')
print("Presentation created successfully: XYZ_Mobile_App_Diagnostic.pptx")
