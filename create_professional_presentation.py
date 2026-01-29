"""
Professional Consulting-Grade Presentation Generator
Modern-minimal design with subtle yellow + red accents
Executive client proposal style
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from datetime import datetime

# Professional Color Palette - Muted, Elegant
DEEP_RED = RGBColor(180, 50, 50)        # Sophisticated burgundy-red
SOFT_RED = RGBColor(200, 100, 100)      # Muted red for accents
WARM_YELLOW = RGBColor(220, 180, 100)   # Muted gold/yellow
LIGHT_YELLOW = RGBColor(250, 240, 210)  # Cream/pale yellow
CREAM = RGBColor(252, 248, 240)         # Warm off-white
IVORY = RGBColor(255, 253, 248)         # Clean ivory background
CHARCOAL = RGBColor(60, 60, 60)         # Deep gray text
STEEL = RGBColor(100, 100, 100)         # Secondary text
LIGHT_GRAY = RGBColor(200, 200, 200)    # Subtle lines
SOFT_WHITE = RGBColor(245, 245, 245)    # Section backgrounds
ACCENT_RED = RGBColor(170, 60, 60)      # Accent line color

def add_clean_header(slide, height=0.9):
    """Minimal header with subtle accent line"""
    # Thin top accent line
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.03))
    accent.fill.solid()
    accent.fill.fore_color.rgb = WARM_YELLOW
    accent.line.fill.background()
    
    # Subtle header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(height))
    header.fill.solid()
    header.fill.fore_color.rgb = IVORY
    header.line.color.rgb = LIGHT_GRAY
    header.line.width = Pt(0.5)
    return header

def add_clean_footer(slide):
    """Minimal footer with thin separator"""
    # Separator line
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.1), Inches(9), Inches(0.015))
    sep.fill.solid()
    sep.fill.fore_color.rgb = LIGHT_GRAY
    sep.line.fill.background()
    
    # Footer text area
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.15), Inches(9), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "XYZ Company | Performance Diagnostic Proposal | January 2026"
    p.font.size = Pt(8)
    p.font.color.rgb = STEEL
    p.font.name = "Calibri Light"
    return footer

def add_section_number(slide, number, left=0.5, top=0.25):
    """Add elegant section number"""
    num_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(0.6), Inches(0.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WARM_YELLOW
    return num_box

def add_slide_title(slide, title, left=0.5, top=0.35, font_size=24, color=CHARCOAL):
    """Clean, professional slide title"""
    title_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Calibri Light"
    return title_box

def add_subtitle(slide, subtitle, left=0.5, top=1.0, color=STEEL):
    """Professional subtitle"""
    sub_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(9), Inches(0.4))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = color
    p.font.name = "Calibri Light"
    return sub_box

def add_body_bullets(slide, bullets, left=0.5, top=1.4, width=9, font_size=11, line_spacing=14):
    """Clean bullet points with proper hierarchy"""
    body_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = body_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = CHARCOAL
        p.font.name = "Calibri"
        p.space_before = Pt(0)
        p.space_after = Pt(line_spacing)
        p.level = 0
    return body_box

def add_horizontal_line(slide, top, left=0.5, width=9, color=LIGHT_GRAY):
    """Subtle horizontal divider"""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line

def add_accent_line(slide, top, left=0.5, width=1.5):
    """Yellow accent underline for emphasis"""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = WARM_YELLOW
    line.line.fill.background()
    return line

def add_content_box(slide, content, left, top, width, height, 
                   bg_color=CREAM, border_color=LIGHT_GRAY, font_size=10):
    """Elegant content box with subtle border"""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), 
                                  Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = border_color
    box.line.width = Pt(0.5)
    
    text_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.1), 
                                         Inches(width - 0.3), Inches(height - 0.2))
    tf = text_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = CHARCOAL
        p.font.name = "Calibri"
        p.space_after = Pt(4)
    return box

def add_numbered_item(slide, number, title, description, left, top, width=4.2):
    """Professional numbered list item"""
    # Number
    num_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(0.3), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(number) + "."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DEEP_RED
    p.font.name = "Calibri Light"
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(left + 0.35), Inches(top), Inches(width - 0.4), Inches(0.3))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CHARCOAL
    p.font.name = "Calibri"
    
    # Description
    if description:
        desc_box = slide.shapes.add_textbox(Inches(left + 0.35), Inches(top + 0.25), 
                                            Inches(width - 0.4), Inches(0.6))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = description
        p.font.size = Pt(9)
        p.font.color.rgb = STEEL
        p.font.name = "Calibri Light"

def create_professional_slide(prs, title, bullets, subtitle=""):
    """Standard professional content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_clean_header(slide)
    add_clean_footer(slide)
    
    add_slide_title(slide, title)
    if subtitle:
        add_subtitle(slide, subtitle)
        add_body_bullets(slide, bullets, top=1.5)
    else:
        add_body_bullets(slide, bullets, top=1.3)
    return slide

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ============ SLIDE 1 — ELEGANT COVER ============
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

# Clean white background
bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = IVORY
bg.line.fill.background()

# Thin top line
line_top = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.015))
line_top.fill.solid()
line_top.fill.fore_color.rgb = LIGHT_GRAY
line_top.line.fill.background()

# Subtle accent bar
accent_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.25), Inches(2), Inches(0.04))
accent_bar.fill.solid()
accent_bar.fill.fore_color.rgb = WARM_YELLOW
accent_bar.line.fill.background()

# Main title
title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(1.2))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "XYZ Mobile App"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = CHARCOAL
p.font.name = "Calibri Light"
p.alignment = PP_ALIGN.LEFT

p = tf.add_paragraph()
p.text = "Performance & Release Management Diagnostic"
p.font.size = Pt(26)
p.font.bold = False
p.font.color.rgb = STEEL
p.font.name = "Calibri Light"
p.alignment = PP_ALIGN.LEFT
p.space_before = Pt(8)

# Bottom line
line_bottom = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.8), Inches(8.4), Inches(0.015))
line_bottom.fill.solid()
line_bottom.fill.fore_color.rgb = LIGHT_GRAY
line_bottom.line.fill.background()

# Company and date
meta_box = slide1.shapes.add_textbox(Inches(0.8), Inches(6.1), Inches(8.4), Inches(1.0))
tf = meta_box.text_frame
p = tf.paragraphs[0]
p.text = "Prepared for XYZ Company"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = DEEP_RED
p.font.name = "Calibri"
p.alignment = PP_ALIGN.LEFT

p = tf.add_paragraph()
p.text = "January 2026"
p.font.size = Pt(11)
p.font.color.rgb = STEEL
p.font.name = "Calibri Light"
p.alignment = PP_ALIGN.LEFT

# Small accent element
small_accent = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(6.5), Inches(0.7), Inches(0.04))
small_accent.fill.solid()
small_accent.fill.fore_color.rgb = WARM_YELLOW
small_accent.line.fill.background()

# ============ SLIDE 2 — UNDERSTANDING SCOPE ============
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_clean_header(slide2)
add_clean_footer(slide2)

add_slide_title(slide2, "Our Understanding of Scope")
add_accent_line(slide2, 0.75)
add_subtitle(slide2, "Client seeks a diagnostic-driven assessment to address the following objectives:")

objectives = [
    "Analyze mobile app latency across Home, Insurance, Spend Track, Quiz and other critical user flows",
    "Identify root causes behind extended load times (6 seconds observed vs. 2–3 second market benchmark)",
    "Understand app size inflation patterns and growth trajectory (Android: 160MB to 400+MB installed; iOS: 402MB)",
    "Determine feasibility and requirements for transitioning to a monthly release cycle",
    "Recommend specific fixes backed by measurable Root Cause Analysis—no assumptions",
    "Provide a North Star performance vision to guide long-term optimization strategy"
]
add_body_bullets(slide2, objectives, top=1.6)

# ============ SLIDE 3 — SCOPE OF DIAGNOSTIC ============
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_clean_header(slide3)
add_clean_footer(slide3)

add_slide_title(slide3, "Scope of Diagnostic")
add_accent_line(slide3, 0.75)

# Left column - AB Team
add_subtitle(slide3, "AB Team Resources & Activities", left=0.5, top=1.1)

add_numbered_item(slide3, 1, "Mobile Performance Lead (Flutter)", "Senior specialist in mobile performance optimization", 0.5, 1.6)
add_numbered_item(slide3, 2, "Mobile Performance Engineer", "Technical implementation and profiling expertise", 0.5, 2.5)
add_numbered_item(slide3, 3, "Engineering Manager", "Program management and client coordination", 0.5, 3.4)

add_horizontal_line(slide3, 4.3, left=0.5, width=4.2)

activities = [
    "Access setup: source code repository, UAT builds, deployment environments",
    "Comprehensive journey mapping and technical architecture walkthrough",
    "Environment validation and build readiness confirmation"
]
add_body_bullets(slide3, activities, left=0.5, top=4.4, width=4.2, font_size=10)

# Vertical divider
vdiv = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(1.3), Inches(0.015), Inches(5.2))
vdiv.fill.solid()
vdiv.fill.fore_color.rgb = LIGHT_GRAY
vdiv.line.fill.background()

# Right column - Client Inputs
add_subtitle(slide3, "Required Client Inputs", left=5.3, top=1.1)

inputs = [
    "Latest production build packages (APK for Android, IPA for iOS)",
    "Access credentials for Analytics dashboards and CMS platforms",
    "Complete API documentation and endpoint specifications",
    "Release pipeline documentation and deployment procedures",
    "Comprehensive third-party SDK inventory and integration details"
]
add_body_bullets(slide3, inputs, left=5.3, top=1.6, width=4.2, font_size=10)

# ============ SLIDE 4 — NORTH STAR VISION ============
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_clean_header(slide4)
add_clean_footer(slide4)

add_slide_title(slide4, "North Star Vision")
add_accent_line(slide4, 0.75)
add_subtitle(slide4, "Performance benchmarks aligned with industry standards")

# Benchmark table headers
headers = ["Metric", "Target", "Current", "Gap"]
x_positions = [0.5, 3.0, 5.0, 7.5]
col_widths = [2.3, 1.8, 1.8, 1.8]

# Header row
for i, (header, x, w) in enumerate(zip(headers, x_positions, col_widths)):
    cell = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.6), Inches(w), Inches(0.45))
    cell.fill.solid()
    cell.fill.fore_color.rgb = CREAM
    cell.line.color.rgb = LIGHT_GRAY
    cell.line.width = Pt(0.5)
    
    text = slide4.shapes.add_textbox(Inches(x), Inches(1.7), Inches(w), Inches(0.35))
    tf = text.text_frame
    p = tf.paragraphs[0]
    p.text = header
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = CHARCOAL
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

# Data rows
benchmarks = [
    ("Primary screen load time", "2 seconds", "6 seconds", "67% reduction"),
    ("Tab-switch latency", "150–250 ms", ">500 ms", "50%+ improvement"),
    ("App size (installed)", "30–40% reduction", "Growing", "Size optimization"),
    ("API latency (critical)", "<150 ms", "300+ ms", "Sub-150ms target"),
    ("Frame render stability", "<16 ms/frame", "Janky frames", "Smooth UX"),
    ("Release cadence", "Monthly train", "Ad-hoc", "Predictable cycle")
]

y_pos = 2.1
for metric, target, current, gap in benchmarks:
    # Row background
    row_bg = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.55))
    row_bg.fill.solid()
    row_bg.fill.fore_color.rgb = IVORY
    row_bg.line.color.rgb = LIGHT_GRAY
    row_bg.line.width = Pt(0.25)
    
    # Metric
    m_text = slide4.shapes.add_textbox(Inches(0.6), Inches(y_pos + 0.15), Inches(2.2), Inches(0.35))
    tf = m_text.text_frame
    p = tf.paragraphs[0]
    p.text = metric
    p.font.size = Pt(9)
    p.font.color.rgb = CHARCOAL
    p.font.name = "Calibri"
    
    # Target
    t_text = slide4.shapes.add_textbox(Inches(3.0), Inches(y_pos + 0.15), Inches(1.8), Inches(0.35))
    tf = t_text.text_frame
    p = tf.paragraphs[0]
    p.text = target
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(60, 120, 60)
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    
    # Current
    c_text = slide4.shapes.add_textbox(Inches(5.0), Inches(y_pos + 0.15), Inches(1.8), Inches(0.35))
    tf = c_text.text_frame
    p = tf.paragraphs[0]
    p.text = current
    p.font.size = Pt(9)
    p.font.color.rgb = DEEP_RED
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    
    # Gap
    g_text = slide4.shapes.add_textbox(Inches(7.5), Inches(y_pos + 0.15), Inches(1.8), Inches(0.35))
    tf = g_text.text_frame
    p = tf.paragraphs[0]
    p.text = gap
    p.font.size = Pt(9)
    p.font.color.rgb = STEEL
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    
    y_pos += 0.58

# Note
note_box = slide4.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(9), Inches(0.4))
tf = note_box.text_frame
p = tf.paragraphs[0]
p.text = "Note: These are reference benchmarks only. Actual commitments will be determined following completion of Root Cause Analysis."
p.font.size = Pt(9)
p.font.italic = True
p.font.color.rgb = STEEL
p.font.name = "Calibri Light"

# ============ SLIDE 5 — ASSUMPTIONS ============
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_clean_header(slide5)
add_clean_footer(slide5)

add_slide_title(slide5, "Key Assumptions")
add_accent_line(slide5, 0.75)

assumptions = [
    ("Client Access", "All required access (source code, builds, dashboards) will be provided by the client in a timely manner to maintain project schedule."),
    ("Technical Constraints", "Third-party SDK behavior and CMS platform limitations may restrict the scope of possible optimizations."),
    ("Scope Boundaries", "No changes to backend systems or CMS platforms unless explicitly included within the engagement scope."),
    ("RCA-Driven Outcomes", "Root Cause Analysis outcomes will determine feasibility—not all identified issues may be technically fixable."),
    ("Data-Driven Approach", "All recommendations will be measurable and derived from systematic profiling and telemetry data analysis."),
    ("UX Stability", "Business-driven UI/UX changes are considered out of scope unless mutually agreed upon via change request process."),
    ("Advisory Role", "Release Management recommendations are advisory; implementation requires client DevOps team involvement.")
]

y_pos = 1.4
for i, (title, desc) in enumerate(assumptions):
    x = 0.5 if i % 2 == 0 else 5.2
    if i > 0 and i % 2 == 0:
        y_pos += 1.5
    
    # Number circle
    circle = slide5.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y_pos), Inches(0.25), Inches(0.25))
    circle.fill.solid()
    circle.fill.fore_color.rgb = WARM_YELLOW
    circle.line.fill.background()
    
    # Number
    num_text = slide5.shapes.add_textbox(Inches(x), Inches(y_pos + 0.02), Inches(0.25), Inches(0.25))
    tf = num_text.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = IVORY
    p.alignment = PP_ALIGN.CENTER
    
    # Title
    title_box = slide5.shapes.add_textbox(Inches(x + 0.35), Inches(y_pos), Inches(4.0), Inches(0.3))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = CHARCOAL
    p.font.name = "Calibri"
    
    # Description
    desc_box = slide5.shapes.add_textbox(Inches(x + 0.35), Inches(y_pos + 0.28), Inches(4.0), Inches(0.9))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(9)
    p.font.color.rgb = STEEL
    p.font.name = "Calibri Light"

# ============ SLIDE 6 — ARCHITECTURE ============
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_clean_header(slide6)
add_clean_footer(slide6)

add_slide_title(slide6, "Architecture & Design Considerations")
add_accent_line(slide6, 0.75)

considerations = [
    "Modular, layered architecture assessment with dependency mapping",
    "API sequencing optimization and dependency chain analysis",
    "Asynchronous vs synchronous rendering path evaluation",
    "Third-party SDK footprint analysis and load behavior profiling",
    "Lazy-loading feasibility study for non-critical resources",
    "Asset compression strategies and intelligent caching mechanisms",
    "Separation of concerns for future scalability and maintainability",
    "Release governance framework and branching strategy review"
]
add_body_bullets(slide6, considerations, top=1.4)

# ============ SLIDE 7 — DIAGNOSTIC VIEW ============
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_clean_header(slide7)
add_clean_footer(slide7)

add_slide_title(slide7, "Proposed Diagnostic Architecture Review")
add_accent_line(slide7, 0.75)

review_areas = [
    "Mobile application frontend architecture (Flutter framework deep-dive)",
    "API Gateway interaction patterns and latency bottlenecks",
    "Content Management System driven modules and payload optimization",
    "Third-party SDK integrations (Fly, MarTech, Payments, Firebase)",
    "Performance telemetry flows and monitoring infrastructure",
    "Release pipeline automation and CI/CD workflow assessment"
]
add_body_bullets(slide7, review_areas, top=1.4)

# ============ SLIDE 8 — TIMELINE ============
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_clean_header(slide8)
add_clean_footer(slide8)

add_slide_title(slide8, "Diagnostic Approach")
add_accent_line(slide8, 0.75)
add_subtitle(slide8, "Week 0 through Week 4 execution plan")

# Timeline visualization
phases = [
    ("Week 0", "Setup & Access", ["Access provisioning", "Environment configuration", "Build validation"]),
    ("Week 1", "Profiling & Analysis", ["Load-time benchmarking", "API call mapping", "Size breakdown", "Cache profiling"]),
    ("Week 2", "RCA & Observations", ["Root-cause identification", "Fixable vs dependency analysis", "SDK constraints", "Pipeline review"]),
    ("Week 3", "Validation", ["Joint walkthrough", "Client team validation", "Feasibility confirmation", "Release alignment"]),
    ("Week 4", "Final Report", ["Target alignment", "Uplift estimation", "Fixes list", "Execution plan"])
]

# Timeline line
timeline = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(8.4), Inches(0.02))
timeline.fill.solid()
timeline.fill.fore_color.rgb = LIGHT_GRAY
timeline.line.fill.background()

# Phase columns
x_start = 1.0
col_width = 1.6
for i, (week, phase, items) in enumerate(phases):
    x = x_start + (i * col_width)
    
    # Phase header
    header = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.5), Inches(1.4), Inches(0.5))
    header.fill.solid()
    header.fill.fore_color.rgb = CREAM if i % 2 == 0 else LIGHT_YELLOW
    header.line.color.rgb = WARM_YELLOW
    header.line.width = Pt(1)
    
    week_text = slide8.shapes.add_textbox(Inches(x), Inches(1.58), Inches(1.4), Inches(0.4))
    tf = week_text.text_frame
    p = tf.paragraphs[0]
    p.text = week
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = CHARCOAL
    p.alignment = PP_ALIGN.CENTER
    
    # Phase name
    phase_box = slide8.shapes.add_textbox(Inches(x), Inches(2.1), Inches(1.4), Inches(0.6))
    tf = phase_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = phase
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = DEEP_RED
    p.alignment = PP_ALIGN.CENTER
    
    # Items
    items_box = slide8.shapes.add_textbox(Inches(x + 0.1), Inches(2.8), Inches(1.3), Inches(3.0))
    tf = items_box.text_frame
    tf.word_wrap = True
    for item in items:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(8)
        p.font.color.rgb = STEEL
        p.space_after = Pt(3)

# ============ SLIDE 9 — TEAM ============
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_clean_header(slide9)
add_clean_footer(slide9)

add_slide_title(slide9, "Team Structure")
add_accent_line(slide9, 0.75)

# Team table headers
headers = ["Role", "Count", "Responsibility"]
x_positions = [0.5, 4.2, 5.0]
col_widths = [3.6, 0.7, 4.0]

for i, (header, x, w) in enumerate(zip(headers, x_positions, col_widths)):
    cell = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.5), Inches(w), Inches(0.45))
    cell.fill.solid()
    cell.fill.fore_color.rgb = CREAM
    cell.line.color.rgb = LIGHT_GRAY
    cell.line.width = Pt(0.5)
    
    text = slide9.shapes.add_textbox(Inches(x), Inches(1.6), Inches(w), Inches(0.35))
    tf = text.text_frame
    p = tf.paragraphs[0]
    p.text = header
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = CHARCOAL
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

# Team data
team = [
    ("Mobile Performance Lead", "1", "Performance profiling, rendering optimization, app load analysis"),
    ("Mobile Engineer", "2", "Code analysis, architecture assessment, technical implementation"),
    ("Solution Architect", "1", "Engineering management, program oversight, strategic architecture")
]

y_pos = 2.0
for role, count, resp in team:
    # Role cell
    role_bg = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y_pos), Inches(3.6), Inches(0.6))
    role_bg.fill.solid()
    role_bg.fill.fore_color.rgb = IVORY
    role_bg.line.color.rgb = LIGHT_GRAY
    role_bg.line.width = Pt(0.25)
    
    role_text = slide9.shapes.add_textbox(Inches(0.6), Inches(y_pos + 0.18), Inches(3.4), Inches(0.3))
    tf = role_text.text_frame
    p = tf.paragraphs[0]
    p.text = role
    p.font.size = Pt(10)
    p.font.color.rgb = CHARCOAL
    
    # Count cell
    count_bg = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.2), Inches(y_pos), Inches(0.7), Inches(0.6))
    count_bg.fill.solid()
    count_bg.fill.fore_color.rgb = LIGHT_YELLOW
    count_bg.line.color.rgb = LIGHT_GRAY
    count_bg.line.width = Pt(0.25)
    
    count_text = slide9.shapes.add_textbox(Inches(4.2), Inches(y_pos + 0.18), Inches(0.7), Inches(0.3))
    tf = count_text.text_frame
    p = tf.paragraphs[0]
    p.text = count
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = CHARCOAL
    p.alignment = PP_ALIGN.CENTER
    
    # Responsibility cell
    resp_bg = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(y_pos), Inches(4.0), Inches(0.6))
    resp_bg.fill.solid()
    resp_bg.fill.fore_color.rgb = IVORY
    resp_bg.line.color.rgb = LIGHT_GRAY
    resp_bg.line.width = Pt(0.25)
    
    resp_text = slide9.shapes.add_textbox(Inches(5.1), Inches(y_pos + 0.12), Inches(3.8), Inches(0.4))
    tf = resp_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = resp
    p.font.size = Pt(9)
    p.font.color.rgb = STEEL
    
    y_pos += 0.65

# Total
total_box = slide9.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.4))
tf = total_box.text_frame
p = tf.paragraphs[0]
p.text = "Total Team: 3 Members (Diagnostic Phase)"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = DEEP_RED
p.font.name = "Calibri"

# ============ SLIDE 10 — GANTT ============
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_clean_header(slide10)
add_clean_footer(slide10)

add_slide_title(slide10, "Project Timeline")
add_accent_line(slide10, 0.75)
add_subtitle(slide10, "Week-level Gantt view")

# Gantt header
weeks = ["Week 0", "Week 1", "Week 2", "Week 3", "Week 4"]
for i, week in enumerate(weeks):
    header = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                       Inches(4.0 + i * 1.15), Inches(1.5),
                                       Inches(1.1), Inches(0.4))
    header.fill.solid()
    header.fill.fore_color.rgb = CREAM
    header.line.color.rgb = LIGHT_GRAY
    header.line.width = Pt(0.5)
    
    text = slide10.shapes.add_textbox(Inches(4.0 + i * 1.15), Inches(1.58), Inches(1.1), Inches(0.3))
    tf = text.text_frame
    p = tf.paragraphs[0]
    p.text = week
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = CHARCOAL
    p.alignment = PP_ALIGN.CENTER

# Gantt bars
tasks = [
    ("Setup & Access", 0, 1),
    ("Profiling & Analysis", 1, 1),
    ("RCA & Observations", 2, 1),
    ("Discussions & Validation", 3, 1),
    ("Final Report", 4, 1)
]

y_pos = 2.2
for task, start, duration in tasks:
    # Task label
    label = slide10.shapes.add_textbox(Inches(0.4), Inches(y_pos + 0.1), Inches(3.4), Inches(0.35))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = task
    p.font.size = Pt(10)
    p.font.color.rgb = CHARCOAL
    
    # Bar
    bar = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(4.0 + start * 1.15), Inches(y_pos),
                                    Inches(duration * 1.05), Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = WARM_YELLOW
    bar.line.color.rgb = ACCENT_RED
    bar.line.width = Pt(1)
    
    y_pos += 0.55

# ============ SLIDE 11 — COMMERCIAL ============
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_clean_header(slide11)
add_clean_footer(slide11)

add_slide_title(slide11, "Commercial Structure")
add_accent_line(slide11, 0.75)

# Left box - Diagnostic
left_box = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0.5), Inches(1.5),
                                     Inches(4.4), Inches(3.5))
left_box.fill.solid()
left_box.fill.fore_color.rgb = CREAM
left_box.line.color.rgb = LIGHT_GRAY
left_box.line.width = Pt(1)

left_header = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.5), Inches(1.5),
                                        Inches(4.4), Inches(0.6))
left_header.fill.solid()
left_header.fill.fore_color.rgb = WARM_YELLOW
left_header.line.fill.background()

left_title = slide11.shapes.add_textbox(Inches(0.5), Inches(1.62), Inches(4.4), Inches(0.4))
tf = left_title.text_frame
p = tf.paragraphs[0]
p.text = "Diagnostic Phase"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = CHARCOAL
p.alignment = PP_ALIGN.CENTER

diag_items = [
    "Duration: 4 weeks",
    "Pricing: Fixed fee structure",
    "Resources: 4–5 FTEs (1 month equivalent)",
    "Deliverables: Analysis, RCA, reporting, advisory"
]
y = 2.25
for item in diag_items:
    item_box = slide11.shapes.add_textbox(Inches(0.7), Inches(y), Inches(4.0), Inches(0.35))
    tf = item_box.text_frame
    p = tf.paragraphs[0]
    p.text = "• " + item
    p.font.size = Pt(10)
    p.font.color.rgb = CHARCOAL
    y += 0.4

# Right box - Execution
right_box = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(5.2), Inches(1.5),
                                      Inches(4.4), Inches(3.5))
right_box.fill.solid()
right_box.fill.fore_color.rgb = IVORY
right_box.line.color.rgb = LIGHT_GRAY
right_box.line.width = Pt(1)

right_header = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         Inches(5.2), Inches(1.5),
                                         Inches(4.4), Inches(0.6))
right_header.fill.solid()
right_header.fill.fore_color.rgb = LIGHT_YELLOW
right_header.line.fill.background()

right_title = slide11.shapes.add_textbox(Inches(5.2), Inches(1.62), Inches(4.4), Inches(0.4))
tf = right_title.text_frame
p = tf.paragraphs[0]
p.text = "Execution Phase"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = CHARCOAL
p.alignment = PP_ALIGN.CENTER

exec_items = [
    "Pricing: Estimated post-diagnostic",
    "Scope: Dependent on fixable items identified",
    "Timeline: To be determined",
    "Resources: Flexible allocation"
]
y = 2.25
for item in exec_items:
    item_box = slide11.shapes.add_textbox(Inches(5.4), Inches(y), Inches(4.0), Inches(0.35))
    tf = item_box.text_frame
    p = tf.paragraphs[0]
    p.text = "• " + item
    p.font.size = Pt(10)
    p.font.color.rgb = CHARCOAL
    y += 0.4

# Arrow
arrow = slide11.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                  Inches(4.7), Inches(2.8),
                                  Inches(0.7), Inches(0.25))
arrow.fill.solid()
arrow.fill.fore_color.rgb = WARM_YELLOW
arrow.line.color.rgb = ACCENT_RED
arrow.line.width = Pt(1)

# ============ SLIDE 12 — RISKS ============
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_clean_header(slide12)
add_clean_footer(slide12)

add_slide_title(slide12, "Risks & Dependencies")
add_accent_line(slide12, 0.75)

risks = [
    "Third-party SDK limitations may constrain optimization options",
    "CMS payload restrictions could impact content delivery performance",
    "Launch-time API dependencies may block user experience improvements",
    "Device fragmentation and low-RAM behavior variations across Android devices",
    "Release process maturity level may affect implementation timeline",
    "Environment availability and access provisioning delays"
]
add_body_bullets(slide12, risks, top=1.4)

# ============ SLIDE 13 — FINAL OUTCOME ============
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
add_clean_header(slide13)
add_clean_footer(slide13)

add_slide_title(slide13, "Final Deliverables")
add_accent_line(slide13, 0.75)
add_subtitle(slide13, "A North Star aligned, data-backed, feasible performance roadmap")

deliverables = [
    "Detailed assessment of what can be improved with specific technical recommendations",
    "Clear documentation of what cannot be improved and underlying technical constraints",
    "Quantified expected performance uplift ranges per optimization category",
    "Recommended team structure and implementation timeline for execution phase",
    "Monthly release train readiness assessment with specific capability gaps identified"
]
add_body_bullets(slide13, deliverables, top=1.6)

# Bottom highlight box
highlight = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(0.5), Inches(6.0),
                                      Inches(9), Inches(0.6))
highlight.fill.solid()
highlight.fill.fore_color.rgb = CREAM
highlight.line.color.rgb = WARM_YELLOW
highlight.line.width = Pt(2)

highlight_text = slide13.shapes.add_textbox(Inches(0.5), Inches(6.15), Inches(9), Inches(0.4))
tf = highlight_text.text_frame
p = tf.paragraphs[0]
p.text = "North Star Aligned • Data-Backed • Feasible & Actionable"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = DEEP_RED
p.alignment = PP_ALIGN.CENTER

# Save with timestamp
timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
output_name = f'XYZ_Mobile_App_Diagnostic_Professional_{timestamp}.pptx'
output_path = f'/Users/Mudassar.Hakim/tempfiles/{output_name}'

prs.save(output_path)
print(f"✓ Professional presentation created: {output_name}")
print(f"  Style: Consulting-grade, modern-minimal, subtle yellow/red accents")
